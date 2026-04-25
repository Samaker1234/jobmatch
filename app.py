from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for
from flask_mail import Mail, Message
from authlib.integrations.flask_client import OAuth
from dotenv import load_dotenv
import os
import re
import PyPDF2
import logging
import shutil
import random

# Load environment variables from .env file
load_dotenv()
from collections import Counter
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import io
import warnings
import json
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from ai_job_generator import JobOfferAIGenerator
from cv_improver import improve_cv, CVImprover
from models import db, User, CVAnalysis, JobProfile, Notification, JobOffer, ResetCode, Candidacy
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from gemini_service import gemini_service
from emails import send_code_email
from forgot import forgot_bp
import secrets
import psutil
import resend

warnings.filterwarnings('ignore')

# Télécharger les ressources NLTK nécessaires
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')

from config import config as app_config

# Choix de configuration via la variable d'environnement FLASK_CONFIG (development/production/testing)
env = os.environ.get('FLASK_CONFIG', 'development')
app = Flask(__name__)

# --- FIX EHLO HOSTNAME WITH COMMA ---
import smtplib
orig_smtp_init = smtplib.SMTP.__init__
def patched_smtp_init(self, *args, **kwargs):
    if 'local_hostname' not in kwargs or kwargs['local_hostname'] is None:
        kwargs['local_hostname'] = 'localhost'
    orig_smtp_init(self, *args, **kwargs)
smtplib.SMTP.__init__ = patched_smtp_init
# ------------------------------------

app.config.from_object(app_config.get(env, app_config['default']))

# Garantit une SECRET_KEY si non fournie
app.config.setdefault('SECRET_KEY', secrets.token_hex(16))
app.config.setdefault('MAX_CONTENT_LENGTH', 50 * 1024 * 1024)  # 50MB max
app.config['JSON_AS_ASCII'] = False  # Support UTF-8 characters (emojis) in JSON responses

# Initialiser la base de données
db.init_app(app)

# Initialiser Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login_page'

# Initialiser Flask-Mail
mail = Mail(app)

# Initialiser OAuth (Google)
oauth = OAuth(app)
google = oauth.register(
    name='google',
    client_id=app.config.get('GOOGLE_CLIENT_ID') or os.environ.get('GOOGLE_CLIENT_ID'),
    client_secret=app.config.get('GOOGLE_CLIENT_SECRET') or os.environ.get('GOOGLE_CLIENT_SECRET'),
    access_token_url='https://oauth2.googleapis.com/token',
    authorize_url='https://accounts.google.com/o/oauth2/v2/auth',
    server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
    client_kwargs={
        'scope': 'openid email profile',
        'prompt': 'select_account'
    }
)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# Créer les tables au démarrage
with app.app_context():
    db.create_all()



# Configurer le logging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# Configurer Resend
resend.api_key = os.environ.get('RESEND_API_KEY')

# Enregistrer les Blueprints
app.register_blueprint(forgot_bp)

@app.before_request
def check_maintenance_mode():
    """Vérifier si le site est en maintenance"""
    # Autoriser l'accès aux admins et à la page superadmin
    if request.path.startswith('/static') or request.path.startswith('/api/admin'):
        return None
    
    if request.path == '/superadmin' or request.path == '/connexion' or request.path == '/logout':
        return None

    # Check maintenance mode in DB
    with app.app_context():
        try:
            is_maint = SystemSetting.get_val('maintenance_mode', '0') == '1'
            if is_maint and not (current_user.is_authenticated and current_user.is_admin):
                return render_template('maintenance.html'), 503
        except:
            pass
    return None

# ============= FONCTIONS UTILITAIRES =============

def extract_text_from_pdf(pdf_file):
    """Extrait le texte d'un fichier PDF avec fallback IA pour les scans"""
    try:
        # 1. Tentative d'extraction classique (rapide)
        pdf_file.seek(0)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        
        # 2. Si l'extraction classique échoue ou donne peu de texte, utiliser Gemini (OCR intelligent)
        if not text or len(text.strip()) < 50:
            print(f">>> DEBUG: Texte insuffisant ({len(text) if text else 0} chars) via PyPDF2. Tentative Gemini...")
            logger.warning(f"ℹ Texte insuffisant ({len(text) if text else 0} chars) via PyPDF2, tentative avec Gemini AI...")
            pdf_file.seek(0)
            pdf_bytes = pdf_file.read()
            ai_text = gemini_service.extract_text_from_pdf_ai(pdf_bytes)
            if ai_text:
                print(">>> DEBUG: Gemini a réussi l'extraction!")
                return ai_text
            else:
                print(">>> DEBUG: Gemini a retourné un texte vide.")
                logger.error("✗ Échec de l'extraction via Gemini AI (réponse vide ou erreur)")
        
        return text
    except Exception as e:
        logger.error(f"Erreur extraction PDF: {e}")
        # Tentative finale via IA même si PyPDF2 a crashé
        try:
            pdf_file.seek(0)
            pdf_bytes = pdf_file.read()
            return gemini_service.extract_text_from_pdf_ai(pdf_bytes)
        except:
            return None

def clean_text(text):
    """Nettoie et normalise le texte"""
    text = text.lower()
    text = re.sub(r'[^a-zéèêëàâäùûüôöçîï0-9\s]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def remove_stopwords(text):
    """Supprime les mots vides"""
    stop_words = set(stopwords.words('french'))
    tokens = word_tokenize(text)
    return ' '.join([word for word in tokens if word not in stop_words and len(word) > 2])

def calculate_similarity(cv_text, job_text):
    """Calcule la similarité cosinus entre le CV et l'offre"""
    vectorizer = TfidfVectorizer(max_features=500, analyzer='char', ngram_range=(2, 3))
    try:
        tfidf_matrix = vectorizer.fit_transform([cv_text, job_text])
        similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
        return float(similarity) * 100
    except:
        return 0

def extract_keywords(text, min_length=3):
    """Extrait les mots-clés pertinents du texte"""
    text = clean_text(text)
    text = remove_stopwords(text)
    tokens = word_tokenize(text)
    filtered_tokens = [token for token in tokens if len(token) > min_length]
    counter = Counter(filtered_tokens)
    return counter.most_common(20)

def find_missing_keywords(cv_keywords, job_keywords):
    """Trouve les mots-clés de l'offre absents du CV"""
    cv_dict = dict(cv_keywords)
    job_dict = dict(job_keywords)
    
    missing = {}
    for keyword, count in job_keywords:
        if keyword not in cv_dict:
            missing[keyword] = count
    
    return sorted(missing.items(), key=lambda x: x[1], reverse=True)

def calculate_detailed_score(cv_text, job_text, similarity_score):
    """Calcule un score détaillé en analysant plusieurs facteurs"""
    cv_words = set(clean_text(cv_text).split())
    job_words = set(clean_text(job_text).split())
    
    if len(job_words) > 0:
        keyword_coverage = len(cv_words & job_words) / len(job_words) * 100
    else:
        keyword_coverage = 0
    
    # Calculate similarity weighted by keyword density
    final_score = (similarity_score * 0.5 + keyword_coverage * 0.5)
    
    # Bonus for specific key matches (optional)
    if final_score > 0 and final_score < 100:
        # Prevent scores from being too low for good matches
        final_score = min(100, final_score * 1.1)
    
    return {
        'similarity': similarity_score,
        'coverage': keyword_coverage,
        'final': final_score
    }

def get_score_color(score):
    """Retourne la classe CSS en fonction du score"""
    if score >= 70:
        return "score-high"
    elif score >= 50:
        return "score-medium"
    else:
        return "score-low"

def extract_cv_details(text):
    """Extrait les informations structurées d'un CV pour le template premium"""
    details = {
        'name': 'Candidat',
        'first_name': 'Prénom',
        'last_name': 'NOM',
        'job_title': 'Expert IT',
        'contact': {'location': '', 'email': '', 'phone': '', 'portfolio': ''},
        'languages': [],
        'skills': [],
        'experiences': [],
        'education': []
    }
    
    # Extraction nom (Très basique, prend la 1ère ligne significative)
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if lines:
        full_name = lines[0]
        details['name'] = full_name
        parts = full_name.split()
        if len(parts) >= 2:
            details['first_name'] = parts[0]
            details['last_name'] = " ".join(parts[1:])
    
    # Emails & Téléphones
    emails = re.findall(r'[a-z0-9\.\-+_]+@[a-z0-9\.\-+_]+\.[a-z]+', text.lower())
    if emails: details['contact']['email'] = emails[0]
    
    phones = re.findall(r'(?:(?:\+|00)33|0)\s*[1-9](?:[\s.-]*\d{2}){4}', text)
    if phones: details['contact']['phone'] = phones[0]
    
    # Job Title (cherche des mots clés après le nom)
    job_titles = ["développeur", "analyste", "ingénieur", "consultant", "manager", "data scientist"]
    for line in lines[1:5]:
        for title in job_titles:
            if title in line.lower():
                details['job_title'] = line.upper()
                break
    
    # Fake splitting for sections (Production real NLP would be better)
    # This is a fallback-rich version for the UX
    details['skills'] = ["Python", "Flask", "React", "SQL", "Git", "Docker"] # Fallback logic
    
    # Logic simplified for demo purposes: in a real app, uses regex for "Expérience", "Formation" etc.
    details['experiences'] = [
        {
            'title': "PROJET D'ANALYSE INTELLIGENTE",
            'date': "2024 - Présent",
            'company': "JobMatch Platform",
            'details': ["Développement de l'interface premium", "Optimisation des scores de similarité"]
        }
    ]
    
    details['education'] = [
        {
            'degree': "FORMATION SUPÉRIEURE EN INFORMATIQUE",
            'date': "En cours",
            'school': "Université Labé"
        }
    ]
    
    return details

# ============= ROUTES =============

@app.route('/loading')
def loading():
    """Page de chargement"""
    return render_template('loading.html')

@app.route('/')
def index():
    """Redirige vers le tableau de bord automatiquement"""
    return redirect('/connexion')

@app.route('/dashboard')
def dashboard():
    """Page du tableau de bord (rendu côté client)"""
    return render_template('dashboard.html', active_page='dashboard')

@app.route('/cv-analysis/<int:analysis_id>')
@login_required
def cv_analysis_detail(analysis_id):
    """Page affichant les détails d'une analyse spécifique"""
    from models import CVAnalysis
    import json
    
    analysis = CVAnalysis.query.get_or_404(analysis_id)
    
    # Security check
    if analysis.user_id != current_user.id and not current_user.is_admin:
        flash("Vous n'êtes pas autorisé à voir cette analyse.", "danger")
        return redirect(url_for('dashboard'))
        
    try:
        missing_keywords = json.loads(analysis.missing_keywords) if analysis.missing_keywords else []
    except:
        missing_keywords = []
        
    return render_template('cv_analysis.html', analysis=analysis, missing_keywords=missing_keywords, active_page='dashboard')

@app.route('/api/dashboard-data')
def dashboard_data():
    """Fournit les données du tableau de bord au format JSON"""
    from models import CVAnalysis
    from datetime import datetime, timedelta
    from flask import jsonify # Ensure jsonify is imported
    
    if current_user.is_authenticated:
        analyses = CVAnalysis.query.filter_by(user_id=current_user.id).order_by(CVAnalysis.created_at.desc()).all()
    else:
        analyses = []
    
    total_analyses = len(analyses)
    avg_score = sum(a.final_score for a in analyses) / total_analyses if total_analyses > 0 else 0
    best_score = max(a.final_score for a in analyses) if total_analyses > 0 else 0
    
    improvement = 0
    if total_analyses >= 2:
        last_score = analyses[0].final_score
        prev_score = analyses[1].final_score
        if prev_score > 0:
            improvement = ((last_score - prev_score) / prev_score) * 100

    dist = {'high': 0, 'medium': 0, 'low': 0}
    for a in analyses:
        if a.final_score >= 80: dist['high'] += 1
        elif a.final_score >= 50: dist['medium'] += 1
        else: dist['low'] += 1
    
    weekly_labels = []
    weekly_values = []
    days_fr = ["Lun", "Mar", "Mer", "Jeu", "Ven", "Sam", "Dim"]
    
    for i in range(6, -1, -1):
        date = datetime.utcnow().date() - timedelta(days=i)
        day_label = days_fr[date.weekday()]
        weekly_labels.append(day_label)
        day_analyses = [a.final_score for a in analyses if a.created_at.date() == date]
        weekly_values.append(max(day_analyses) if day_analyses else 0)
    
    if total_analyses == 0:
        weekly_labels = ["Lun", "Mar", "Mer", "Jeu", "Ven", "Sam", "Dim"]
        weekly_values = [0, 0, 0, 0, 0, 0, 0]

    latest_analyses_data = []
    for a in analyses[:5]:
        latest_analyses_data.append({
            'id': a.id,
            'job_title': a.job_title,
            'filename': a.filename,
            'created_at': a.created_at.strftime('%d/%m/%Y'),
            'final_score': round(a.final_score)
        })

    user_info = {
        'is_authenticated': current_user.is_authenticated,
        'firstname': current_user.firstname if current_user.is_authenticated else 'Alice'
    }

    return jsonify({
        'user': user_info,
        'total_analyses': total_analyses,
        'avg_score': round(avg_score, 1),
        'best_score': round(best_score, 1),
        'improvement': round(improvement, 2),
        'dist': dist,
        'weekly_labels': weekly_labels,
        'weekly_values': weekly_values,
        'latest_analyses': latest_analyses_data
    })

@app.route('/job-generator')
@login_required
def job_generator():
    """Page du générateur d'offres IA"""
    return render_template('job_generator_v3.html', user=current_user, active_page='job-generator')

@app.route('/about')
def about_page():
    """Page À propos - Accessible à tous"""
    return render_template('about.html', user=current_user if current_user.is_authenticated else None, active_page='about')

@app.route('/contact')
def contact_page():
    """Page Contact"""
    return render_template('contact.html', user=current_user, active_page='contact')

@app.route('/faq')
def faq_page():
    """Page FAQ"""
    return render_template('faq.html', user=current_user, active_page='faq')

@app.route('/legal')
def legal_page():
    """Page Légal / Mentions légales"""
    return render_template('legal.html', user=current_user, active_page='legal')

@app.route('/cv-improver')
@login_required
def cv_improver():
    """Page d'amélioration de CV"""
    return render_template('cv_improver_v3.html', user=current_user, active_page='cv-improver')

@app.route('/my-offers')
@login_required
def my_offers():
    """Page des offres sauvegardées"""
    return render_template('my_offers_v3.html', user=current_user, active_page='my-offers')

@app.route('/settings', methods=['GET', 'POST'])
@login_required
def settings():
    """Page des paramètres avec fonctionnalités avancées"""
    if request.method == 'POST':
        action = request.form.get('action')
        
        try:
            # 1. Mise à jour informations générales
            if action == 'update_profile':
                firstname = request.form.get('firstname')
                lastname = request.form.get('lastname')
                profession = request.form.get('profession')
                
                if firstname and lastname:
                    current_user.firstname = firstname
                    current_user.lastname = lastname
                    current_user.profession = profession
                    db.session.commit()
                    # Flash success (optionnel)
            
            # 2. Gestion des images
            elif action == 'update_images':
                from werkzeug.utils import secure_filename
                
                # Dossier d'upload (Assurez-vous qu'il existe)
                upload_folder = os.path.join(app.static_folder, 'uploads', 'users')
                os.makedirs(upload_folder, exist_ok=True)
                
                # Profile Picture
                if 'profile_picture' in request.files:
                    file = request.files['profile_picture']
                    if file and file.filename != '':
                        filename = secure_filename(f"{current_user.id}_profile_{int(datetime.utcnow().timestamp())}.jpg") # Force .jpg for simplicity or keep ext
                        file.save(os.path.join(upload_folder, filename))
                        current_user.profile_picture = url_for('static', filename=f'uploads/users/{filename}')
                
                # Banner
                if 'banner_image' in request.files:
                    file = request.files['banner_image']
                    if file and file.filename != '':
                        filename = secure_filename(f"{current_user.id}_banner_{int(datetime.utcnow().timestamp())}.jpg")
                        file.save(os.path.join(upload_folder, filename))
                        current_user.banner_image = url_for('static', filename=f'uploads/users/{filename}')
                
                db.session.commit()

            # 3. Sécurité (Mot de passe)
            elif action == 'change_password':
                current_pw = request.form.get('current_password')
                new_pw = request.form.get('new_password')
                confirm_pw = request.form.get('confirm_password')
                
                if not current_user.check_password(current_pw):
                    # Flash error: Mot de passe actuel incorrect
                    pass 
                elif new_pw != confirm_pw:
                    # Flash error: Les mots de passe ne correspondent pas
                    pass
                else:
                    current_user.set_password(new_pw)
                    db.session.commit()
                    # Flash success
            
            # 4. Suppression de compte
            elif action == 'delete_account':
                # Nettoyage des fichiers (optionnel mais recommandé)
                db.session.delete(current_user)
                db.session.commit()
                logout_user()
                return redirect(url_for('login_page'))

        except Exception as e:
            db.session.rollback()
            logger.error(f"Erreur settings ({action}): {str(e)}")
            
        return redirect(url_for('settings'))

    return render_template('settings.html', user=current_user, active_page='settings')

@app.route('/alerts')
@login_required
def alerts():
    """Page des alertes emploi"""
    return render_template('alerts.html', user=current_user, active_page='alerts')

@app.route('/subscriptions')
@login_required
def subscriptions():
    """Page des abonnements"""
    return render_template('subscriptions_v3.html', user=current_user, active_page='subscriptions')



@app.route('/superadmin')
@login_required
def superadmin():
    """Page Super Admin (Static Template, Data fetched via JS)"""
    if not current_user.is_admin:
        return redirect('/dashboard')
    
    return render_template('superadmin.html', user=current_user, active_page='superadmin')

@app.route('/api/admin/system-insights')
@login_required
def admin_system_insights():
    """Analyse les stats système via Gemini pour donner des conseils à l'admin"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
    
    try:
        # Collecter les métriques
        user_count = User.query.count()
        analysis_count = CVAnalysis.query.count()
        cpu = psutil.cpu_percent()
        ram = psutil.virtual_memory().percent
        
        stats_summary = f"""
        Utilisateurs totaux: {user_count}
        Analyses de CV effectuées: {analysis_count}
        Charge CPU: {cpu}%
        Utilisation RAM: {ram}%
        """
        
        prompt = f"""
        En tant qu'assistant d'administration système intelligent, analysez ces métriques de la plateforme JobMatch :
        {stats_summary}
        
        Fournissez 3 conseils stratégiques courts pour l'administrateur (croissance, maintenance ou performance).
        Répondez en français au format JSON : 
        {{
            "insights": ["conseil 1", "conseil 2", "conseil 3"],
            "status": "Excellent/Bon/Attention"
        }}
        """
        
        if gemini_service.enabled:
            import google.generativeai as genai
            model = genai.GenerativeModel('gemini-pro')
            response = model.generate_content(prompt)
            insights = json.loads(response.text.replace('```json', '').replace('```', '').strip())
            return jsonify(insights)
        else:
            return jsonify({
                'insights': ["Activez Gemini pour des conseils IA.", "Maintenez le système à jour.", "Surveillez les logs."],
                'status': "Stable"
            })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/logs')
@login_required
def get_admin_logs():
    """Récupère les dernières lignes du log système"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
    
    # Simulons ou lisons un fichier log si nécessaire
    # Pour l'instant, on renvoie des événements récents de la DB
    logs = [
        {"time": datetime.now().strftime("%H:%M:%S"), "event": "Nouvelle connexion admin", "level": "INFO"},
        {"time": datetime.now().strftime("%H:%M:%S"), "event": "Backup automatique réussi", "level": "SUCCESS"},
        {"time": datetime.now().strftime("%H:%M:%S"), "event": "Analyse CV optimisée (Gemini)", "level": "AI"}
    ]
    return jsonify({'logs': logs})

@app.route('/api/admin/users/<int:user_id>', methods=['DELETE'])
@login_required
def delete_user(user_id):
    """Supprimer un utilisateur (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        user = User.query.get(user_id)
        if not user:
            return jsonify({'error': 'Utilisateur non trouvé'}), 404
            
        # Bloquer uniquement l'auto-suppression
        if user.id == current_user.id:
            return jsonify({'error': 'Vous ne pouvez pas supprimer votre propre compte'}), 400
            
        db.session.delete(user)
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Utilisateur supprimé avec succès'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users', methods=['POST'])
@login_required
def create_user_admin():
    """Créer un utilisateur (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        data = request.get_json()
        firstname = data.get('firstname', '').strip()
        lastname = data.get('lastname', '').strip()
        email = data.get('email', '').strip()
        password = data.get('password', '').strip()
        is_admin = data.get('is_admin', False)
        
        if not all([firstname, lastname, email, password]):
            return jsonify({'error': 'Tous les champs sont requis'}), 400
            
        if User.query.filter_by(email=email).first():
            return jsonify({'error': 'Cet email existe déjà'}), 400
            
        new_user = User(
            firstname=firstname,
            lastname=lastname,
            email=email,
            profession='Administrateur' if is_admin else 'Membre',
            is_admin=is_admin
        )
        new_user.set_password(password)
        
        db.session.add(new_user)
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Utilisateur créé avec succès'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/users/<int:user_id>', methods=['GET'])
@login_required
def get_user_admin(user_id):
    """Récupérer un utilisateur (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    user = User.query.get(user_id)
    if not user:
        return jsonify({'error': 'Utilisateur non trouvé'}), 404
        
    return jsonify({
        'id': user.id,
        'firstname': user.firstname,
        'lastname': user.lastname,
        'email': user.email,
        'is_admin': user.is_admin
    })

@app.route('/api/admin/users/<int:user_id>', methods=['PUT'])
@login_required
def update_user_admin(user_id):
    """Modifier un utilisateur (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        user = User.query.get(user_id)
        if not user:
            return jsonify({'error': 'Utilisateur non trouvé'}), 404
            
        data = request.get_json()
        
        if 'firstname' in data:
            user.firstname = data['firstname'].strip()
        if 'lastname' in data:
            user.lastname = data['lastname'].strip()
        if 'email' in data:
            email = data['email'].strip()
            if email != user.email:
                if User.query.filter_by(email=email).first():
                    return jsonify({'error': 'Cet email existe déjà'}), 400
                user.email = email
        if 'password' in data and data['password'].strip():
            user.set_password(data['password'].strip())
        
        # Changement de rôle (is_admin)
        if 'is_admin' in data:
            # Sécurité : ne pas pouvoir se retirer ses propres droits admin
            if user.id == current_user.id and data['is_admin'] is False:
                pass # Ignorer silencieusement ou retourner une erreur si on veut être strict
            else:
                user.is_admin = bool(data['is_admin'])
                # Mettre à jour la profession pour refleter le rôle
                if user.is_admin:
                    user.profession = 'Administrateur'
                else:
                    user.profession = 'Membre'
            
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Utilisateur mis à jour avec succès'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/clear-cache', methods=['POST'])
@login_required
def clear_cache():
    """Vider le cache système (Simulé pour l'instant)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
    
    # In a real app, you would clear Redis or filesystem cache here
    # For now, we simulate it
    try:
        # Simulation of cache clearing
        import time
        time.sleep(1) 
        return jsonify({'success': True, 'message': 'Cache système vidé avec succès'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/settings', methods=['GET', 'POST'])
@login_required
def admin_settings():
    """Gérer les paramètres système (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    if request.method == 'GET':
        settings = SystemSetting.query.all()
        return jsonify({s.key: s.value for s in settings})
        
    try:
        data = request.get_json()
        for key, value in data.items():
            setting = SystemSetting.query.filter_by(key=key).first()
            if not setting:
                setting = SystemSetting(key=key, value=str(value))
                db.session.add(setting)
            else:
                setting.value = str(value)
        
        db.session.commit()
        return jsonify({'success': True, 'message': 'Paramètres mis à jour'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/dashboard-data', methods=['GET'])
@login_required
def admin_dashboard_data():
    """Récupérer toutes les données pour le dashboard superadmin"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    from models import User, CVAnalysis, JobOffer
    from sqlalchemy import func
    from datetime import datetime, timedelta
    import calendar

    # Statistiques de base
    stats = {
        'user_count': User.query.count(),
        'offer_count': JobOffer.query.count(),
        'analysis_count': CVAnalysis.query.count(),
        'active_today': User.query.filter(User.created_at >= datetime.utcnow().date()).count()
    }
    
    # Weekly Data
    weekly = {'labels': [], 'users': [], 'analyses': []}
    days_fr = ["Lun", "Mar", "Mer", "Jeu", "Ven", "Sam", "Dim"]
    for i in range(6, -1, -1):
        date = datetime.utcnow().date() - timedelta(days=i)
        weekly['labels'].append(days_fr[date.weekday()])
        weekly['users'].append(User.query.filter(func.date(User.created_at) == date).count())
        weekly['analyses'].append(CVAnalysis.query.filter(func.date(CVAnalysis.created_at) == date).count())
    
    # Monthly Data
    monthly = {'labels': [], 'users': [], 'analyses': []}
    for i in range(5, -1, -1):
        m_date = datetime.utcnow().replace(day=1) - timedelta(days=i*30)
        m_idx, y_val = m_date.month, m_date.year
        monthly['labels'].append(calendar.month_name[m_idx][:3])
        monthly['users'].append(User.query.filter(func.extract('month', User.created_at) == m_idx, func.extract('year', User.created_at) == y_val).count())
        monthly['analyses'].append(CVAnalysis.query.filter(func.extract('month', CVAnalysis.created_at) == m_idx, func.extract('year', CVAnalysis.created_at) == y_val).count())

    user_dist = {
        'admins': User.query.filter_by(is_admin=True).count(),
        'users': User.query.filter_by(is_admin=False).count()
    }

    # Latest users for the table
    latest_users = User.query.order_by(User.created_at.desc()).limit(15).all()
    user_list = [{
        'id': u.id,
        'firstname': u.firstname,
        'lastname': u.lastname,
        'email': u.email,
        'is_admin': u.is_admin,
        'created_at': u.created_at.strftime('%d %b %Y') if u.created_at else 'N/A'
    } for u in latest_users]

    return jsonify({
        'stats': stats,
        'weekly': weekly,
        'monthly': monthly,
        'user_distribution': user_dist,
        'users': user_list,
        'current_user_email': current_user.email
    })

@app.route('/api/admin/broadcast', methods=['POST'])
@login_required
def admin_broadcast():
    """Diffuser une annonce à tous les utilisateurs (Superadmin uniquement)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        data = request.get_json()
        title = data.get('title', 'Annonce Système')
        message = data.get('message')
        notif_type = data.get('type', 'info')
        
        if not message:
            return jsonify({'error': 'Message requis'}), 400
            
        users = User.query.all()
        
        # 1. Créer des notifications internes
        for user in users:
            notif = Notification(
                user_id=user.id,
                title=title,
                message=message,
                type=notif_type
            )
            db.session.add(notif)
            
        # 2. Envoyer un email récapitulatif (optionnel, via Resend)
        # Note: on réutilise la logique de email-users si besoin
        
        db.session.commit()
        return jsonify({'success': True, 'message': f'Annonce diffusée à {len(users)} utilisateurs'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/admin/email-users', methods=['POST'])
@login_required
def email_all_users():
    """Envoyer un email à tous les utilisateurs (Simulé)"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403

    try:
        users = User.query.all()
        emails = [u.email for u in users if u.email]

        if not emails:
            return jsonify({'success': True, 'message': 'Aucun utilisateur à contacter'})

        # Send broadcast email
        # Note: If Resend domain is not verified, you can only send to the verified email
        # We send as bcc to hide users from each other
        params = {
            "from": "JobMatch <samakedelamou858@gmail.com>",
            "to": [emails[0]],  # Send to one as 'to'
            "bcc": emails[1:] if len(emails) > 1 else [],
            "subject": "Message de l'Administration JobMatch",
            "html": f"<strong>Bonjour à tous,</strong><br><br>Ceci est un message important de l'administration JobMatch Console.<br><br>Cordialement,<br>L'équipe Admin",
        }

        r = resend.Emails.send(params)

        return jsonify({'success': True, 'message': f'Email envoyé avec succès à {len(emails)} utilisateurs via Resend'})
    except Exception as e:
        logger.error(f"Erreur Resend: {str(e)}")
        return jsonify({'error': f'Erreur lors de l\'envoi des emails : {str(e)}'}), 500


@app.route('/api/admin/test-email', methods=['POST'])
@login_required
def admin_test_email():
    """Envoie un email de test via Flask-Mail pour vérifier la configuration SMTP."""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403

    try:
        data = request.get_json(silent=True) or {}
        # Priorité : email fourni dans la requête, sinon email de l'utilisateur, sinon MAIL_USERNAME
        to_email = data.get('email') or getattr(current_user, 'email', None) or app.config.get('MAIL_USERNAME')

        if not to_email:
            return jsonify({'error': 'Aucune adresse email de destination disponible'}), 400

        test_code = str(data.get('code', '123456'))

        # Import local pour éviter les effets de bord au démarrage
        from emails import send_code_email

        ok = send_code_email(to_email, test_code)
        if not ok:
            return jsonify({
                'error': 'Échec de l\'envoi du mail de test. Vérifiez la configuration SMTP et les logs serveur.'
            }), 500

        logger.info(f"Email de test Flask-Mail envoyé à {to_email}")
        return jsonify({'success': True, 'message': f'Email de test envoyé à {to_email}'}), 200
    except Exception as e:
        logger.error(f'Erreur lors de l\'envoi de l\'email de test: {e}', exc_info=True)
        return jsonify({'error': str(e)}), 500


@app.route('/api/admin/backup-db', methods=['POST'])
@login_required
def backup_database():
    """Sauvegarder la base de données"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        import shutil
        import datetime
        
        # Source database path (adjust if needed)
        db_path = os.path.join(app.instance_path, 'jobmatch.db')
        
        if not os.path.exists(db_path):
            return jsonify({'error': 'Fichier de base de données introuvable'}), 404
            
        # Backup destination
        backup_dir = os.path.join(app.instance_path, 'backups')
        os.makedirs(backup_dir, exist_ok=True)
        
        timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
        backup_filename = f'jobmatch_backup_{timestamp}.db'
        backup_path = os.path.join(backup_dir, backup_filename)
        
        shutil.copy2(db_path, backup_path)
        
        return jsonify({'success': True, 'message': f'Sauvegarde créée : {backup_filename}'})
    except Exception as e:
        return jsonify({'error': f'Erreur de sauvegarde : {str(e)}'}), 500

@app.route('/api/admin/health', methods=['GET'])
@login_required
def system_health_check():
    """Vérification complète de la santé du système"""
    if not current_user.is_admin:
        return jsonify({'error': 'Non autorisé'}), 403
        
    try:
        # System Metrics
        cpu_usage = psutil.cpu_percent(interval=None)
        memory = psutil.virtual_memory()
        disk = shutil.disk_usage("/")
        
        # Database Check
        try:
            db.session.execute(db.text('SELECT 1'))
            db_status = "Connected"
        except:
            db_status = "Disconnected"
            
        health_data = {
            'status': 'Healthy' if db_status == "Connected" and cpu_usage < 90 else 'Warning',
            'cpu': cpu_usage,
            'memory': {
                'total': round(memory.total / (1024**3), 2),
                'used': round(memory.used / (1024**3), 2),
                'percent': memory.percent
            },
            'disk': {
                'total': round(disk.total / (1024**3), 2),
                'free': round(disk.free / (1024**3), 2),
                'percent': round((disk.used / disk.total) * 100, 1)
            },
            'database': db_status,
            'timestamp': datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        }
        
        return jsonify(health_data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health')
def health():
    """Endpoint de vérification de santé"""
    return jsonify({
        'status': 'healthy',
        'service': 'JobMatch API',
        'version': '1.0.0'
    }), 200

@app.route('/download-cv-template')
def download_cv_template():
    """Télécharger le modèle de CV"""
    try:
        pdf_path = os.path.join(os.path.dirname(__file__), 'CV_Exemple.pdf')
        return send_file(
            pdf_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='CV_Modele_JobMatch.pdf'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/analyze', methods=['POST'])
def analyze():
    """Endpoint API pour l'analyse enrichie avec Gemini"""
    try:
        # Récupérer le fichier PDF et l'offre d'emploi
        if 'pdf_file' not in request.files:
            logger.warning("Analyse échouée: Aucun fichier pdf_file dans request.files")
            return jsonify({'error': 'Aucun fichier PDF fourni'}), 400
        
        pdf_file = request.files['pdf_file']
        job_offer = request.form.get('job_offer', '').strip()
        threshold = int(request.form.get('threshold', 3))
        
        if pdf_file.filename == '':
            logger.warning("Analyse échouée: Nom de fichier vide")
            return jsonify({'error': 'Aucun fichier sélectionné'}), 400
        
        if not job_offer or len(job_offer) < 20:
            logger.warning(f"Analyse échouée: Offre d'emploi trop courte ({len(job_offer) if job_offer else 0} chars)")
            return jsonify({'error': 'Offre d\'emploi trop courte (minimum 20 caractères)'}), 400
        
        # Extraction du texte du CV
        cv_text = extract_text_from_pdf(pdf_file)
        
        if not cv_text or len(cv_text.strip()) == 0:
            logger.warning("Analyse échouée: Texte extrait du PDF est vide")
            return jsonify({'error': 'Impossible d\'extraire le texte du PDF'}), 400
        
        print(f">>> DEBUG: Texte extrait ({len(cv_text)} chars). Validation CV...")
            
        # NOUVEAU: Vérifier si c'est vraiment un CV
        cv_check = gemini_service.is_cv_valid(cv_text)
        if not cv_check.get('valid'):
            logger.warning(f"Analyse échouée: Document non valide selon Gemini - {cv_check.get('reason')}")
            return jsonify({
                'error': 'Document non valide',
                'details': cv_check.get('reason')
            }), 400
        
        print(">>> DEBUG: CV validé. Calcul des scores...")
        
        # Nettoyage des textes
        cv_clean = clean_text(cv_text)
        job_clean = clean_text(job_offer)
        
        # Calcul du score de similarité
        similarity_score = calculate_similarity(cv_clean, job_clean)
        
        # Score détaillé
        scores = calculate_detailed_score(cv_clean, job_clean, similarity_score)
        
        # Extraction des mots-clés
        cv_keywords = extract_keywords(cv_clean)
        job_keywords = extract_keywords(job_clean)
        
        # Mots-clés manquants
        missing_keywords = find_missing_keywords(cv_keywords, job_keywords)
        missing_keywords = [kw for kw in missing_keywords if kw[1] >= threshold][:15]
        
        print(f">>> DEBUG: Score de base calculé: {scores['final']}. Enrichissement IA...")
        
        # Recommandations basiques
        recommendation = ""
        if scores['final'] >= 80:
            recommendation = "✅ **Match Exceptionnel!** Votre CV est parfaitement aligné avec les exigences majeures de cette offre. Vos chances de sélection sont très élevées."
        elif scores['final'] >= 65:
            recommendation = "🌟 **Très Bon Profil** - Votre CV correspond bien à l'offre. Quelques ajustements mineurs sur les mots-clés suggérés pourraient optimiser votre score."
        elif scores['final'] >= 45:
            recommendation = "⚠️ **Potentiel Intéressant** - Vous possédez les bases, mais votre CV manque de termes techniques cruciaux listés ci-après pour passer les systèmes de filtrage."
        else:
            recommendation = "ℹ️ **Optimisation Nécessaire** - Votre profil semble s'écarter des attentes de l'offre. Nous vous conseillons de réviser vos expériences pour mettre en avant les compétences demandées."
        
        # Enrichissement avec Gemini AI si activé
        gemini_analysis = {}
        if gemini_service.enabled:
            gemini_analysis = gemini_service.analyze_cv_compatibility(cv_text[:2000], job_offer[:1500])
        
        print(">>> DEBUG: Analyse Gemini terminée. Préparation de la réponse...")
        
        # Déterminer le score final (priorité à Gemini)
        final_score = scores['final']
        recommendation_final = recommendation
        matching_skills_list = [kw for kw, _ in cv_keywords[:10]]
        missing_skills_list = [kw for kw, _ in missing_keywords]
        improvement_plan = "Consultez les recommandations pour optimiser votre score."

        if gemini_analysis and 'error' not in gemini_analysis:
            final_score = gemini_analysis.get('compatibility_score', scores['final'])
            recommendation_final = gemini_analysis.get('insight', recommendation)
            matching_skills_list = gemini_analysis.get('matching_skills', matching_skills_list)
            missing_skills_list = gemini_analysis.get('missing_skills', missing_skills_list)
            improvement_plan = gemini_analysis.get('recommendations', improvement_plan)
            # S'assurer que improvement_plan est une chaîne
            if isinstance(improvement_plan, list):
                improvement_plan = " • ".join(improvement_plan)

        # Sauvegarder dans la base de données si l'utilisateur est connecté
        candidacy_id = None
        if current_user.is_authenticated:
            try:
                # Extraire un titre de poste lisible depuis l'offre d'emploi
                job_title_short = job_offer[:80].split('\n')[0].strip() if job_offer else 'Poste analysé'
                
                analysis = CVAnalysis(
                    user_id=current_user.id,
                    filename=pdf_file.filename,
                    job_title=job_clean[:100] + "..." if len(job_clean) > 100 else job_clean,
                    final_score=final_score,
                    similarity_score=scores['similarity'],
                    coverage_score=scores['coverage'],
                    recommendation=recommendation_final,
                    missing_keywords=json.dumps(missing_skills_list)
                )
                db.session.add(analysis)
                
                # Auto-créer une Candidacy pour tout CV analysé
                try:
                    candidacy = Candidacy(
                        user_id=current_user.id,
                        job_title=job_title_short,
                        company='',
                        filename=pdf_file.filename,
                        cv_score=round(final_score, 1),
                        source='analyzed',
                        status='pending',
                        is_submitted=False
                    )
                    db.session.add(candidacy)
                    db.session.commit()
                    candidacy_id = candidacy.id
                except Exception as e:
                    print(f"Erreur de sauvegarde DB Candidacy: {e}")
                    db.session.rollback()
            except Exception as e:
                print(f"Erreur de sauvegarde DB Analysis: {e}")
                db.session.rollback()

        response_data = {
            'success': True,
            'candidacy_id': candidacy_id,
            'scores': {
                'final': round(final_score, 1),
                'similarity': round(scores['similarity'], 1),
                'coverage': round(scores['coverage'], 1)
            },
            # Support pour job_generator_v3.html
            'final_score': round(final_score, 1),
            'match_percentage': round(final_score, 1), # Utiliser le score final pour la compatibilité
            'skill_match_score': round(scores['coverage'], 1), # Score de couverture des compétences
            'matching_skills': matching_skills_list,
            'missing_skills': missing_skills_list,
            'improvement_plan': improvement_plan,
            
            'score_color': get_score_color(final_score),
            'recommendation': recommendation_final,
            'missing_keywords': [{'keyword': kw, 'count': count} for kw, count in missing_keywords],
            'cv_keywords': [{'keyword': kw, 'count': count} for kw, count in cv_keywords[:10]],
            'job_keywords': [{'keyword': kw, 'count': count} for kw, count in job_keywords[:10]]
        }
        
        # Ajouter l'analyse brute Gemini si disponible pour débug
        if gemini_analysis and 'error' not in gemini_analysis:
            response_data['gemini_analysis'] = gemini_analysis
            response_data['ai_powered'] = True
        
        return jsonify(response_data)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ============= AUTHENTIFICATION ROUTES =============

# Autoriser HTTP pour OAuth en local (Développement uniquement)
os.environ['OAUTHLIB_INSECURE_TRANSPORT'] = '1'

# ============= AUTHENTIFICATION GOOGLE =============
@app.route('/login/google')
def login_google():
    """Rediriger l'utilisateur vers la page d'autorisation Google"""
    # Utiliser l'URI configuré dans .env ou générer dynamiquement
    redirect_uri = app.config.get('GOOGLE_REDIRECT_URI') or url_for('authorize', _external=True)
    logger.info(f"OAuth Redirect URI used: {redirect_uri}")
    return google.authorize_redirect(redirect_uri)

@app.route('/authorize')
def authorize():
    """Gérer le retour de l'authentification Google"""
    try:
        token = google.authorize_access_token()
        user_info = token.get('userinfo')
        if not user_info:
            return redirect(url_for('login_page'))
        
        # Rechercher l'utilisateur par Google ID ou Email
        user = User.query.filter_by(google_id=user_info['sub']).first()
        if not user:
            user = User.query.filter_by(email=user_info['email']).first()
            if user:
                # Lier le compte Google à l'utilisateur existant
                user.google_id = user_info['sub']
                if not user.profile_picture:
                    user.profile_picture = user_info.get('picture')
            else:
                # Créer un nouvel utilisateur
                user = User(
                    firstname=user_info.get('given_name', 'Utilisateur'),
                    lastname=user_info.get('family_name', 'Google'),
                    email=user_info['email'],
                    google_id=user_info['sub'],
                    profile_picture=user_info.get('picture'),
                    profession='Membre'
                )
                db.session.add(user)
        
        db.session.commit()
        login_user(user)
        return redirect(url_for('dashboard'))
    except Exception as e:
        logger.error(f"Erreur OAuth Google : {str(e)}")
        return redirect(url_for('login_page'))

@app.route('/connexion')
def login_page():
    """Page de connexion"""
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    return render_template('connexion.html')

@app.route('/api/login', methods=['POST'])
def api_login():
    """Authentification utilisateur via API"""
    data = request.get_json()
    email = data.get('email', '').strip()
    password = data.get('password', '')
    
    if not email or not password:
        return jsonify({'error': 'Email et mot de passe requis'}), 400
        
    user = User.query.filter_by(email=email).first()
    
    if user and user.check_password(password):
        login_user(user, remember=True)
        return jsonify({'success': True, 'message': 'Connexion réussie', 'redirect': '/dashboard'})
    
    return jsonify({'error': 'Email ou mot de passe incorrect'}), 401

@app.route('/logout')
@login_required
def logout():
    """Déconnexion utilisateur"""
    logout_user()
    return redirect(url_for('login_page'))

@app.route('/register')
def register_page():
    """Page d'inscription"""
    return render_template('register.html')

@app.route('/forgot-password')
def forgot_password_page():
    """Page de récupération de mot de passe"""
    return render_template('forgot_password.html')

@app.route('/api/reset-password', methods=['POST'])
def api_reset_password():
    """Vérifier le code et réinitialiser le mot de passe"""
    data = request.get_json()
    email = data.get('email', '').strip()
    code = data.get('code', '').strip()
    new_password = data.get('password', '').strip()
    
    if not all([email, code, new_password]):
        return jsonify({'error': 'Tous les champs sont requis'}), 400
        
    reset_entry = ResetCode.query.filter_by(email=email, code=code).first()
    
    if not reset_entry:
        return jsonify({'error': 'Code ou email invalide'}), 400
        
    if datetime.utcnow() > reset_entry.expires_at:
        return jsonify({'error': 'Le code a expiré'}), 400
        
    user = User.query.filter_by(email=email).first()
    if not user:
        return jsonify({'error': 'Utilisateur non trouvé'}), 400
        
    try:
        user.set_password(new_password)
        # Supprimer le code utilisé
        db.session.delete(reset_entry)
        db.session.commit()
        return jsonify({'success': True, 'message': 'Mot de passe réinitialisé avec succès'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/profile')
@login_required
def profile_page():
    """Page de profil utilisateur"""
    analyses_count = CVAnalysis.query.filter_by(user_id=current_user.id).count()
    recent_analyses = CVAnalysis.query.filter_by(user_id=current_user.id).order_by(CVAnalysis.created_at.desc()).limit(10).all()
    
    return render_template('profile.html', 
                          user=current_user, 
                          analyses_count=analyses_count,
                          recent_analyses=recent_analyses,
                          active_page='profile')

@app.route('/api/upload-profile-picture', methods=['POST'])
@login_required
def upload_profile_picture():
    """Upload une photo de profil"""
    try:
        from werkzeug.utils import secure_filename
        
        logger.info(f"Tentative d'upload de photo profil pour user {current_user.id}")
        
        if 'file' not in request.files:
            logger.warning("Aucun fichier fourni")
            return jsonify({'error': 'Aucun fichier fourni'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            logger.warning("Nom de fichier vide")
            return jsonify({'error': 'Nom de fichier vide'}), 400
        
        # Vérifier le format
        allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        
        if file_ext not in allowed_extensions:
            logger.warning(f"Format non supporté: {file_ext}")
            return jsonify({'error': 'Format non supporté. Utilisez: PNG, JPG, GIF, WEBP'}), 400
        
        # Créer le dossier s'il n'existe pas
        upload_folder = os.path.join(app.static_folder, 'uploads', 'profiles')
        os.makedirs(upload_folder, exist_ok=True)
        
        # Créer le nom de fichier sécurisé
        filename = f"{current_user.id}_profile_{int(datetime.utcnow().timestamp())}.{file_ext}"
        filepath = os.path.join(upload_folder, filename)
        
        logger.info(f"Sauvegarde du fichier: {filepath}")
        file.save(filepath)
        
        # Mettre à jour la base de données
        current_user.profile_picture = f'/static/uploads/profiles/{filename}'
        db.session.commit()
        
        logger.info(f"✓ Photo de profil uploadée pour user {current_user.id}")
        
        return jsonify({
            'success': True,
            'message': 'Photo de profil mise à jour',
            'profile_picture': current_user.profile_picture
        }), 200
        
    except Exception as e:
        logger.error(f"✗ Erreur upload photo profil: {str(e)}", exc_info=True)
        try:
            db.session.rollback()
        except:
            pass
        return jsonify({'error': f'Erreur serveur: {str(e)}'}), 500

@app.route('/api/upload-banner-image', methods=['POST'])
@login_required
def upload_banner_image():
    """Upload une image de bannière"""
    try:
        from werkzeug.utils import secure_filename
        
        logger.info(f"Tentative d'upload de bannière pour user {current_user.id}")
        
        if 'file' not in request.files:
            logger.warning("Aucun fichier fourni")
            return jsonify({'error': 'Aucun fichier fourni'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            logger.warning("Nom de fichier vide")
            return jsonify({'error': 'Nom de fichier vide'}), 400
        
        # Vérifier le format
        allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        
        if file_ext not in allowed_extensions:
            logger.warning(f"Format non supporté: {file_ext}")
            return jsonify({'error': 'Format non supporté. Utilisez: PNG, JPG, GIF, WEBP'}), 400
        
        # Créer le dossier s'il n'existe pas
        upload_folder = os.path.join(app.static_folder, 'uploads', 'profiles')
        os.makedirs(upload_folder, exist_ok=True)
        
        # Créer le nom de fichier sécurisé
        filename = f"{current_user.id}_banner_{int(datetime.utcnow().timestamp())}.{file_ext}"
        filepath = os.path.join(upload_folder, filename)
        
        logger.info(f"Sauvegarde du fichier: {filepath}")
        file.save(filepath)
        
        # Mettre à jour la base de données
        current_user.banner_image = f'/static/uploads/profiles/{filename}'
        db.session.commit()
        
        logger.info(f"✓ Bannière uploadée pour user {current_user.id}")
        
        return jsonify({
            'success': True,
            'message': 'Bannière mise à jour',
            'banner_image': current_user.banner_image
        }), 200
        
    except Exception as e:
        logger.error(f"✗ Erreur upload bannière: {str(e)}", exc_info=True)
        try:
            db.session.rollback()
        except:
            pass
        return jsonify({'error': f'Erreur serveur: {str(e)}'}), 500

@app.route('/team')
def team_page():
    """Page Présentation de l'équipe - Accessible à tous"""
    return render_template('team.html',
                          user=current_user if current_user.is_authenticated else None,
                          active_page='team')

# ============= API AUTHENTIFICATION =============

@app.route('/api/register', methods=['POST'])
def api_register():
    """API d'inscription"""
    try:
        data = request.get_json()
        firstname = data.get('firstname', '').strip()
        lastname = data.get('lastname', '').strip()
        email = data.get('email', '').strip()
        password = data.get('password', '')
        profession = data.get('profession', '')

        # Validation
        if not all([firstname, lastname, email, password, profession]):
            return jsonify({'error': 'Tous les champs sont requis'}), 400

        if len(firstname) < 2:
            return jsonify({'error': 'Le prénom est trop court'}), 400

        if len(lastname) < 2:
            return jsonify({'error': 'Le nom est trop court'}), 400

        if '@' not in email or '.' not in email.split('@')[1]:
            return jsonify({'error': 'Email invalide'}), 400

        if len(password) < 8:
            return jsonify({'error': 'Le mot de passe doit avoir au moins 8 caractères'}), 400

        # Vérifier si l'email existe déjà
        if User.query.filter_by(email=email).first():
            return jsonify({'error': 'Cet email est déjà utilisé'}), 400

        # Créer le nouvel utilisateur
        new_user = User(
            firstname=firstname,
            lastname=lastname,
            email=email,
            profession=profession
        )
        new_user.set_password(password)
        
        db.session.add(new_user)
        db.session.commit()
        
        return jsonify({
            'message': 'Inscription réussie',
            'user': {
                'firstname': firstname,
                'lastname': lastname,
                'email': email,
                'profession': new_user.profession,
                'bio': new_user.bio,
                'achievements': new_user.achievements,
                'skills': new_user.skills
            },
        }), 201

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/update-profile', methods=['POST'])
@login_required
def api_update_profile():
    """API pour mettre à jour les informations du profil utilisateur"""
    try:
        data = request.get_json()
        
        firstname = data.get('firstname', '').strip()
        lastname = data.get('lastname', '').strip()
        email = data.get('email', '').strip()
        profession = data.get('profession', '').strip()
        bio = data.get('bio', '').strip()
        achievements = data.get('achievements', '').strip()
        skills = data.get('skills', '').strip()
        new_password = data.get('new_password', '').strip()

        if not firstname or not lastname or not email:
            return jsonify({'error': 'Le prénom, le nom et l\'email sont obligatoires'}), 400

        # Vérifier si l'email est déjà utilisé par un autre utilisateur
        existing_user = User.query.filter(User.email == email, User.id != current_user.id).first()
        if existing_user:
            return jsonify({'error': 'Cet email est déjà utilisé par un autre compte'}), 400

        # Mettre à jour les infos
        current_user.firstname = firstname
        current_user.lastname = lastname
        current_user.email = email
        current_user.profession = profession
        current_user.bio = bio
        current_user.achievements = achievements
        current_user.skills = skills

        if new_password:
            if len(new_password) < 6:
                return jsonify({'error': 'Le mot de passe doit faire au moins 6 caractères'}), 400
            current_user.set_password(new_password)

        db.session.commit()
        return jsonify({'success': True, 'message': 'Profil mis à jour'})

    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/extract-job-from-image', methods=['POST'])
@login_required
def extract_job_from_image():
    """API pour extraire le texte d'une image d'offre d'emploi"""
    try:
        if 'image' not in request.files:
            return jsonify({'error': 'Aucune image fournie'}), 400
        
        image_file = request.files['image']
        if image_file.filename == '':
            return jsonify({'error': 'Aucun fichier sélectionné'}), 400
        
        # Read image data
        image_data = image_file.read()
        mime_type = image_file.content_type or 'image/jpeg'
        
        # Call gemini service for OCR
        text = gemini_service.extract_text_from_image(image_data, mime_type)
        
        if not text:
            return jsonify({'error': 'Impossible d\'extraire le texte de l\'image'}), 500
            
        return jsonify({
            'success': True,
            'text': text
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notifications', methods=['GET'])
@login_required
def get_notifications():
    """Récupère les notifications de l'utilisateur"""
    try:
        notifications = Notification.query.filter_by(user_id=current_user.id).order_by(Notification.created_at.desc()).limit(20).all()
        return jsonify({
            'success': True,
            'notifications': [{
                'id': n.id,
                'title': n.title,
                'message': n.message,
                'type': n.type,
                'is_read': n.is_read,
                'created_at': n.created_at.strftime('%Y-%m-%d %H:%M:%S')
            } for n in notifications]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notifications/read/<int:notif_id>', methods=['POST'])
@login_required
def mark_notification_read(notif_id):
    """Marque une notification comme lue"""
    try:
        notif = Notification.query.filter_by(id=notif_id, user_id=current_user.id).first()
        if notif:
            notif.is_read = True
            db.session.commit()
            return jsonify({'success': True})
        return jsonify({'error': 'Notification non trouvée'}), 404
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/notifications/read-all', methods=['POST'])
@login_required
def mark_all_notifications_read():
    """Marque toutes les notifications comme lues"""
    try:
        Notification.query.filter_by(user_id=current_user.id, is_read=False).update({Notification.is_read: True})
        db.session.commit()
        return jsonify({'success': True})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

# ============= NOUVELLES FONCTIONNALITÉS =============

@app.route('/api/export-pdf', methods=['POST'])
def export_pdf():
    """Exporte les résultats d'analyse en PDF"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'Aucune donnée fournie'}), 400
        
        # Créer le PDF en mémoire
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch)
        story = []
        styles = getSampleStyleSheet()
        
        # Style personnalisé
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#004E89'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#FF6B35'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        # Titre
        story.append(Paragraph("Rapport d'Analyse JobMatch", title_style))
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph(f"Date: {datetime.now().strftime('%d/%m/%Y %H:%M')}", styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Scores
        story.append(Paragraph("📊 Scores de Compatibilité", heading_style))
        scores_data = [
            ['Métrique', 'Score'],
            ['Score Global', f"{data.get('scores', {}).get('final', 0)}%"],
            ['Similarité Textuelle', f"{data.get('scores', {}).get('similarity', 0)}%"],
            ['Couverture Mots-clés', f"{data.get('scores', {}).get('coverage', 0)}%"]
        ]
        scores_table = Table(scores_data, colWidths=[4*inch, 2*inch])
        scores_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#004E89')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.grey)
        ]))
        story.append(scores_table)
        story.append(Spacer(1, 0.3*inch))
        
        # Recommandation
        story.append(Paragraph("💡 Recommandation", heading_style))
        recommendation = data.get('recommendation', '').replace('**', '').replace('✅', '').replace('⚠️', '').replace('ℹ️', '')
        story.append(Paragraph(recommendation, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Mots-clés manquants
        missing_keywords = data.get('missing_keywords', [])
        if missing_keywords:
            story.append(Paragraph("🔑 Mots-clés à Ajouter", heading_style))
            keywords_text = ", ".join([kw.get('keyword', '') for kw in missing_keywords[:10]])
            story.append(Paragraph(keywords_text, styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
        
        # Construire le PDF
        doc.build(story)
        buffer.seek(0)
        
        return send_file(
            buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f'JobMatch_Analyse_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/suggest-rephrasing', methods=['POST'])
def suggest_rephrasing():
    """Suggère des reformulations pour améliorer le CV"""
    try:
        data = request.get_json()
        cv_text = data.get('cv_text', '')
        job_offer = data.get('job_offer', '')
        
        if not cv_text or not job_offer:
            return jsonify({'error': 'CV et offre requis'}), 400
        
        # Extraire les compétences de l'offre
        job_keywords = extract_keywords(clean_text(job_offer))
        cv_keywords = extract_keywords(clean_text(cv_text))
        
        # Trouver les compétences manquantes
        missing = find_missing_keywords(cv_keywords, job_keywords)
        
        # Générer des suggestions
        suggestions = []
        for keyword, count in missing[:8]:
            priority = 'high' if count >= 4 else 'medium'
            suggestions.append({
                'keyword': keyword,
                'suggestion': f"Intégrez le terme '{keyword}' dans vos réalisations passées pour démontrer une expérience concrète.",
                'priority': priority
            })
        
        return jsonify({
            'success': True,
            'suggestions': suggestions
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/compare-cvs', methods=['POST'])
def compare_cvs():
    """Compare plusieurs CV avec une offre"""
    try:
        if 'pdf_files' not in request.files:
            return jsonify({'error': 'Aucun fichier PDF fourni'}), 400
        
        pdf_files = request.files.getlist('pdf_files')
        job_offer = request.form.get('job_offer', '').strip()
        
        if len(pdf_files) < 2:
            return jsonify({'error': 'Au moins 2 CV sont requis pour la comparaison'}), 400
        
        if not job_offer or len(job_offer) < 50:
            return jsonify({'error': 'Offre d\'emploi invalide'}), 400
        
        job_clean = clean_text(job_offer)
        results = []
        
        for idx, pdf_file in enumerate(pdf_files):
            cv_text = extract_text_from_pdf(pdf_file)
            if not cv_text:
                continue
                
            cv_clean = clean_text(cv_text)
            similarity = calculate_similarity(cv_clean, job_clean)
            scores = calculate_detailed_score(cv_clean, job_clean, similarity)
            
            results.append({
                'filename': pdf_file.filename,
                'index': idx + 1,
                'scores': {
                    'final': round(scores['final'], 1),
                    'similarity': round(scores['similarity'], 1),
                    'coverage': round(scores['coverage'], 1)
                },
                'score_color': get_score_color(scores['final'])
            })
        
        # Trier par score final
        results.sort(key=lambda x: x['scores']['final'], reverse=True)
        
        return jsonify({
            'success': True,
            'comparison': results,
            'best_cv': results[0] if results else None
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/trends', methods=['GET'])
def get_trends():
    """Analyse les tendances réelles des mots-clés depuis la base de données"""
    try:
        # Récupérer toutes les analyses
        analyses = CVAnalysis.query.all()
        
        all_keywords = []
        for a in analyses:
            if a.missing_keywords:
                all_keywords.extend(json.loads(a.missing_keywords))
        
        # Simuler des tendances (en production, comparer avec les données précédentes)
        kw_counts = Counter(all_keywords).most_common(8)
        trends_list = []
        for kw, count in kw_counts:
            trends_list.append({
                'keyword': kw,
                'frequency': count,
                'trend': random.choice(['up', 'stable', 'down'])
            })

        trends = {
            'top_keywords': trends_list,
            'sectors': [
                {'sector': 'Développement Web', 'demand': 45},
                {'sector': 'Data Science', 'demand': 30},
                {'sector': 'DevOps', 'demand': 15},
                {'sector': 'Mobile', 'demand': 10}
            ],
            'avg_score': 62.5,
            'total_analyses': len(analyses)
        }
        
        return jsonify({
            'success': True,
            'trends': trends
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/save-profile', methods=['POST'])
@login_required
def save_profile():
    """Sauvegarde un profil d'offre dans la base de données"""
    try:
        data = request.get_json()
        
        profile = JobProfile(
            user_id=current_user.id,
            title=data.get('title', 'Offre sans titre'),
            company=data.get('company', ''),
            description=data.get('job_offer', ''),
            tags=",".join(data.get('tags', []))
        )
        
        db.session.add(profile)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'profile': {
                'id': profile.id,
                'title': profile.title,
                'company': profile.company,
                'date': profile.created_at.isoformat()
            }
        }), 201
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/dashboard-stats', methods=['GET'])
@login_required
def dashboard_stats():
    """Retourne les statistiques réelles pour le dashboard"""
    try:
        analyses = CVAnalysis.query.filter_by(user_id=current_user.id).order_by(CVAnalysis.created_at.desc()).all()
        
        if not analyses:
            return jsonify({
                'success': True,
                'stats': {
                    'total_analyses': 0,
                    'avg_score': 0,
                    'best_score': 0,
                    'improvement_rate': 0,
                    'scores_distribution': [],
                    'recent_analyses': [],
                    'top_missing_keywords': []
                }
            })

        total = len(analyses)
        avg_score = sum(a.final_score for a in analyses) / total
        best_score = max(a.final_score for a in analyses)
        
        # Calcul du taux d'amélioration (comparaison entre la première et la dernière analyse)
        improvement_rate = 0
        if total >= 2:
            first_score = analyses[-1].final_score
            last_score = analyses[0].final_score
            if first_score > 0:
                improvement_rate = ((last_score - first_score) / first_score) * 100

        # Distribution des scores
        distribution = [
            {'range': '0-30', 'count': len([a for a in analyses if a.final_score <= 30])},
            {'range': '31-50', 'count': len([a for a in analyses if 31 <= a.final_score <= 50])},
            {'range': '51-70', 'count': len([a for a in analyses if 51 <= a.final_score <= 70])},
            {'range': '71-90', 'count': len([a for a in analyses if 71 <= a.final_score <= 90])},
            {'range': '91-100', 'count': len([a for a in analyses if a.final_score > 90])}
        ]

        # Mots-clés manquants les plus fréquents
        all_missing = []
        for a in analyses:
            if a.missing_keywords:
                all_missing.extend(json.loads(a.missing_keywords))
        
        kw_counts = Counter(all_missing).most_common(5)
        top_missing = [{'keyword': kw, 'frequency': count} for kw, count in kw_counts]

        # Récupérer les notifications récentes
        notifications = Notification.query.filter_by(user_id=current_user.id).order_by(Notification.created_at.desc()).limit(3).all()
        
        # Récupérer les offres sauvegardées récentes
        from models import SavedOffer
        saved_offers = SavedOffer.query.filter_by(user_id=current_user.id).order_by(SavedOffer.created_at.desc()).limit(3).all()

        stats = {
            'total_analyses': total,
            'avg_score': round(avg_score, 1),
            'best_score': round(best_score, 1),
            'improvement_rate': round(improvement_rate, 1),
            'scores_distribution': distribution,
            'recent_analyses': [
                {
                    'date': a.created_at.strftime('%d %b %Y'),
                    'score': round(a.final_score, 1),
                    'title': a.job_title or 'Analyse sans titre',
                    'similarity_score': round(a.similarity_score or 0, 1),
                    'coverage_score': round(a.coverage_score or 0, 1)
                } for a in analyses[:5]
            ],
            'top_missing_keywords': top_missing,
            'notifications': [
                {
                    'title': n.title,
                    'message': n.message,
                    'type': n.type,
                    'date': n.created_at.strftime('%Y-%m-%d %H:%M:%S')
                } for n in notifications
            ],
            'saved_offers': [
                {
                    'title': so.offer.title if so.offer else 'Offre supprimée',
                    'company': so.offer.company if so.offer else 'N/A',
                    'status': so.status,
                    'date': so.created_at.strftime('%d %b %Y')
                } for so in saved_offers
            ]
        }
        
        return jsonify({
            'success': True,
            'stats': stats
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate-premium-cv', methods=['POST'])
def generate_premium_cv():
    """Génère un CV PDF premium basé sur le texte du PDF téléversé"""
    try:
        if 'pdf_file' not in request.files:
            return jsonify({'error': 'Aucun fichier PDF fourni'}), 400
        
        pdf_file = request.files['pdf_file']
        cv_text = extract_text_from_pdf(pdf_file)
        
        if not cv_text:
            return jsonify({'error': 'Échec de l\'extraction du texte'}), 400
            
        # Vérification si c'est un CV
        cv_check = gemini_service.is_cv_valid(cv_text)
        if not cv_check.get('valid'):
            return jsonify({'error': 'Document non valide', 'details': cv_check.get('reason')}), 400
            
        # Extraire les infos
        cv_data = extract_cv_details(cv_text)
        
        # Lire le CSS du dossier Format
        css_path = os.path.join(app.root_path, 'Format', 'style.css')
        css_content = ""
        if os.path.exists(css_path):
            with open(css_path, 'r', encoding='utf-8') as f:
                css_content = f.read()
        
        # Rendre le HTML avec Jinja2
        rendered_html = render_template('premium_cv_template.html', 
                                      css_content=css_content,
                                      **cv_data)
        
        # En production, on utiliserait pdfkit ou weasyprint pour HTML -> PDF
        # Ici on simule en renvoyant le HTML pour visualisation ou on génère via reportlab
        # Pour rester dans la stack installée sans installer wkhtmltopdf :
        
        return rendered_html # Retourne le HTML pour que l'user puisse "Imprimer en PDF"
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ============= ROUTES IA - GÉNÉRATEUR D'OFFRES D'EMPLOI =============

@app.route('/api/generate-job-offer', methods=['POST'])
def generate_job_offer():
    """Génère une offre d'emploi via Gemini AI"""
    try:
        if not gemini_service.enabled:
            return jsonify({'error': 'Service Gemini non disponible'}), 503
        
        data = request.get_json()
        job_title = data.get('job_title', '').strip()
        company = data.get('company', '').strip()
        skills = data.get('skills', [])
        experience_level = data.get('experience_level', 'Senior').strip()
        
        # Fallback aux anciens paramètres si nouveaux paramètres non fournis
        if not job_title and data.get('role_type'):
            job_title = data.get('role_type', 'Développeur')
            company = 'Entreprise Technologique'
            skills = ['Python', 'JavaScript', 'SQL']
        
        if not job_title or not company or not skills:
            return jsonify({'error': 'Titre du poste, entreprise et compétences requis'}), 400
        
        # Générer l'offre avec Gemini
        offer_data = gemini_service.generate_job_offers(
            job_title=job_title,
            company=company,
            skills=skills,
            experience_level=experience_level
        )
        
        if 'error' in offer_data:
            return jsonify({'error': 'Erreur lors de la génération de l\'offre', 'details': offer_data.get('error')}), 500
        
        return jsonify({
            'success': True,
            'offer': offer_data
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-test-job-offer', methods=['POST'])
def generate_test_job_offer():
    """Génère une offre d'emploi de test via Gemini pour tester l'optimiseur CV"""
    try:
        import random
        
        # Liste de profils variés pour générer différentes offres
        job_profiles = [
            {
                'title': 'Développeur Full Stack Senior',
                'company': 'TechCorp',
                'skills': ['Python', 'React', 'PostgreSQL', 'Docker', 'AWS'],
                'experience': '5+ ans'
            },
            {
                'title': 'Data Scientist',
                'company': 'DataVision',
                'skills': ['Python', 'Machine Learning', 'TensorFlow', 'SQL', 'Pandas'],
                'experience': '3-5 ans'
            },
            {
                'title': 'DevOps Engineer',
                'company': 'CloudSystems',
                'skills': ['Kubernetes', 'Docker', 'Jenkins', 'Terraform', 'AWS'],
                'experience': '4+ ans'
            },
            {
                'title': 'Frontend Developer',
                'company': 'WebAgency',
                'skills': ['React', 'TypeScript', 'CSS3', 'Redux', 'Next.js'],
                'experience': '2-4 ans'
            },
            {
                'title': 'Product Manager Tech',
                'company': 'InnovateNow',
                'skills': ['Agile', 'Jira', 'Product Strategy', 'UX/UI', 'Analytics'],
                'experience': '5+ ans'
            }
        ]
        
        # Sélectionner un profil aléatoire
        profile = random.choice(job_profiles)
        
        # Utiliser Gemini pour générer une offre complète et réaliste
        if gemini_service.enabled:
            offer_data = gemini_service.generate_job_offers(
                job_title=profile['title'],
                company=profile['company'],
                skills=profile['skills'],
                experience_level=profile['experience']
            )
            
            if "error" in offer_data:
                return jsonify({'error': offer_data['error']}), 500
                
            # Construire le texte de l'offre à partir des données structurées
            job_offer_text = f"Titre: {offer_data.get('title', profile['title'])}\n"
            job_offer_text += f"Entreprise: {offer_data.get('company', profile['company'])}\n\n"
            job_offer_text += f"Description:\n{offer_data.get('description', '')}\n\n"
            job_offer_text += f"Prérequis:\n- " + "\n- ".join(offer_data.get('requirements', profile['skills'])) + "\n\n"
            job_offer_text += f"Avantages:\n- " + "\n- ".join(offer_data.get('benefits', []))
            
            return jsonify({
                'success': True,
                'job_offer': job_offer_text.strip(),
                'profile': profile,
                'is_ai': True
            })
        else:
            # Fallback statique si Gemini est désactivé
            job_offer_text = f"Offre d'emploi fictive pour {profile['title']} chez {profile['company']}.\nCompétences: {', '.join(profile['skills'])}"
            return jsonify({
                'success': True,
                'job_offer': job_offer_text,
                'profile': profile,
                'is_ai': False
            })
            
    except Exception as e:
        logger.error(f"Erreur génération offre test: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-job-offers-batch', methods=['POST'])
def generate_job_offers_batch():
    """Génère un lot d'offres d'emploi"""
    try:
        data = request.get_json()
        count = data.get('count', 10)
        role_type = data.get('role_type', None)
        
        if count < 1 or count > 20: # Limit for Gemini
            count = 10
        
        generator = JobOfferAIGenerator()
        offers = generator.generate_batch(count, role_type)
        
        # SI GEMINI EST ACTIVÉ: Améliorer les descriptions du lot
        if gemini_service.enabled:
            for offer in offers:
                try:
                    # Enrichissement rapide via Gemini
                    prompt = f"Améliorez très brièvement cette description d'offre pour {offer['title']}: {offer['description'][:200]}. Répondez avec la description enrichie uniquement."
                    enriched = gemini_service.model.generate_content(prompt)
                    if enriched and enriched.text:
                        offer['description'] = enriched.text.strip()
                        offer['is_ai_enriched'] = True
                except:
                    pass
        
        return jsonify({
            'success': True,
            'offers': offers,
            'total': len(offers)
        })
    except Exception as e:
        logger.error(f"Erreur batch generation: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate-offers-by-cv', methods=['POST'])
def generate_offers_by_cv():
    """Génère des offres basées sur les compétences du CV"""
    try:
        if 'pdf_file' not in request.files:
            return jsonify({'error': 'Aucun fichier PDF fourni'}), 400
        
        pdf_file = request.files['pdf_file']
        cv_text = extract_text_from_pdf(pdf_file)
        
        if not cv_text:
            return jsonify({'error': 'Échec de l\'extraction du texte'}), 400
            
        # Vérification si c'est un CV
        cv_check = gemini_service.is_cv_valid(cv_text)
        if not cv_check.get('valid'):
            return jsonify({'error': 'Document non valide', 'details': cv_check.get('reason')}), 400
        
        # Extraire les mots-clés du CV
        cv_keywords = extract_keywords(cv_text)
        required_skills = [keyword[0] for keyword in cv_keywords[:10]]
        
        generator = JobOfferAIGenerator()
        
        # Générer plusieurs offres adaptées
        count = request.form.get('count', 5, type=int)
        offers = []
        
        for _ in range(min(count, 5)): # Limited for performance
            offer = generator.generate_by_requirements(required_skills)
            
            # IA Optimization if enabled
            if gemini_service.enabled:
                try:
                    # Optimisation contextuelle
                    prompt = f"Adaptez cette offre pour un profil ayant ces compétences: {', '.join(required_skills)}. Offre: {offer['description'][:200]}. Répondez avec le texte optimisé."
                    response = gemini_service.model.generate_content(prompt)
                    if response and response.text:
                        offer['description'] = response.text.strip()
                        offer['is_personalized_ai'] = True
                except:
                    pass
                    
            # Calculer le score de match avec le CV
            offer['match_score'] = calculate_similarity(cv_text, offer['description'])
            offers.append(offer)
        
        # Trier par score de match
        offers.sort(key=lambda x: x['match_score'], reverse=True)
        
        return jsonify({
            'success': True,
            'offers': offers,
            'cv_skills_detected': required_skills,
            'total': len(offers)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate-specific-offer', methods=['POST'])
def generate_specific_offer():
    """Génère une offre spécifique avec critères personnalisés"""
    try:
        data = request.get_json()
        
        required_skills = data.get('required_skills', [])
        experience_level = data.get('experience_level', None)
        location = data.get('location', None)
        
        if not required_skills:
            return jsonify({'error': 'Compétences requises manquantes'}), 400
        
        generator = JobOfferAIGenerator()
        offer = generator.generate_by_requirements(
            required_skills=required_skills,
            experience_level=experience_level,
            location=location
        )
        
        return jsonify({
            'success': True,
            'offer': offer
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============= ROUTES CV IMPROVER =============

@app.route('/api/improve-cv', methods=['POST'])
def improve_cv_route():
    """Améliore un CV uploadé et retourne le HTML formaté"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'Aucun fichier uploadé'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'Nom de fichier vide'}), 400
        
        # Récupérer l'offre d'emploi si fournie
        job_offer = request.form.get('job_offer', '').strip()
        
        # Extraire le texte du PDF
        cv_text = extract_text_from_pdf(file)
        
        if not cv_text:
            return jsonify({'error': 'Impossible d\'extraire le texte du PDF'}), 400
            
        # MODE TURBO : Optimisation et Structuration en un seul appel
        if gemini_service.enabled:
            # On utilise une seule session Gemini pour TOUT faire (validation + optimisation)
            turbo_result = gemini_service.optimize_cv_and_structure(cv_text, job_offer if job_offer and len(job_offer) >= 50 else None)
            
            if 'error' not in turbo_result:
                # Si non valide, on s'arrête
                if not turbo_result.get('is_valid', True):
                    return jsonify({'error': 'Document non valide', 'details': turbo_result.get('reason_invalid')}), 400
                
                # Récupérer les données structurées et les suggestions
                cv_data = turbo_result.get('cv_data')
                
                # Injecter la photo de profil si disponible
                if current_user.is_authenticated and current_user.profile_picture:
                    cv_data['profile_picture'] = current_user.profile_picture
                
                improver = CVImprover()
                html_res = improver.generate_html_cv(cv_data)
                
                # NOUVEAU: Envoi automatique par mail si l'utilisateur est connecté
                candidacy_id = None
                if current_user.is_authenticated:
                    # 1. Tentative d'envoi d'email
                    try:
                        from emails import send_cv_email
                        pdf_buf = generate_cv_pdf_buffer(cv_data)
                        filename = f"{cv_data.get('nom', 'CV')}_Optimise.pdf"
                        send_cv_email(current_user.email, pdf_buf, filename)
                    except Exception as e:
                        logger.error(f"Erreur envoi auto CV: {str(e)}")
                    
                    # 2. Auto-création de la candidature
                    try:
                        job_title_from_cv = cv_data.get('titre_professionnel', 'Poste optimisé')
                        score_optim = turbo_result.get('score_optimisation', 0)
                        cand = Candidacy(
                            user_id=current_user.id,
                            job_title=job_title_from_cv,
                            company='',
                            filename=file.filename,
                            cv_score=float(score_optim) if score_optim else None,
                            source='optimized',
                            status='pending',
                            is_submitted=False
                        )
                        db.session.add(cand)
                        db.session.commit()
                        candidacy_id = cand.id
                    except Exception as e:
                        logger.error(f"Erreur création candidature auto: {str(e)}")
                        db.session.rollback()
                
                return jsonify({
                    'success': True,
                    'data': cv_data,
                    'improvements': turbo_result.get('suggestions', []),
                    'score_optimisation': turbo_result.get('score_optimisation', 0),
                    'html': html_res,
                    'has_gemini': True,
                    'turbo_mode': True,
                    'candidacy_id': candidacy_id
                })
        
        # Fallback si Gemini désactivé ou erreur
        result = improve_cv(cv_text)
        
        # Envoi auto pour le fallback aussi
        if current_user.is_authenticated:
            try:
                from emails import send_cv_email
                pdf_buf = generate_cv_pdf_buffer(result['data'])
                send_cv_email(current_user.email, pdf_buf, "CV_Optimise.pdf")
            except: pass

        return jsonify({
            'success': True,
            'data': result['data'],
            'improvements': result['improvements'],
            'html': result['html'],
            'has_gemini': False
        })
        
        return jsonify({
            'success': True,
            'data': result['data'],
            'improvements': result['improvements'],
            'html': result['html'],
            'has_gemini': 'gemini_optimization' in result or 'gemini_data' in result
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/download-cv-html', methods=['POST'])
def download_cv_html():
    """Télécharge le CV amélioré en HTML"""
    try:
        data = request.get_json()
        html_content = data.get('html')
        nom = data.get('nom', 'CV')
        
        if not html_content:
            return jsonify({'error': 'Contenu HTML manquant'}), 400
        
        # Créer un fichier HTML
        html_bytes = html_content.encode('utf-8')
        
        return send_file(
            io.BytesIO(html_bytes),
            mimetype='text/html',
            as_attachment=True,
            download_name=f'{nom}_CV_Ameliore.html'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500


def generate_cv_pdf_buffer(cv_data):
    """Génère un buffer PDF à partir des données CV avec un design premium à deux colonnes"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    import io

    buffer = io.BytesIO()
    # Marges minimales car on utilise une table pour le layout
    doc = SimpleDocTemplate(buffer, pagesize=A4, 
                           rightMargin=0, leftMargin=0,
                           topMargin=0, bottomMargin=0)
    
    styles = getSampleStyleSheet()
    
    # --- Couleurs ---
    color_primary = colors.HexColor('#1e293b')
    color_accent = colors.HexColor('#6366f1')
    color_sidebar_bg = colors.HexColor('#f1f5f9')
    color_text = colors.HexColor('#334155')
    color_text_light = colors.HexColor('#64748b')

    # --- Styles de texte ---
    style_name = ParagraphStyle('Name', parent=styles['Normal'], fontSize=28, textColor=color_primary, leading=32, fontName='Helvetica-Bold')
    style_title = ParagraphStyle('JobTitle', parent=styles['Normal'], fontSize=16, textColor=color_text_light, spaceBefore=4, fontName='Helvetica-Bold')
    style_section = ParagraphStyle('Section', parent=styles['Normal'], fontSize=12, textColor=color_primary, fontName='Helvetica-Bold', spaceBefore=15, spaceAfter=8, borderPadding=(0,0,2,0), borderStyle='Solid', borderColor=color_accent)
    style_sidebar_h = ParagraphStyle('SidebarH', parent=styles['Normal'], fontSize=10, textColor=color_primary, fontName='Helvetica-Bold', spaceBefore=15, spaceAfter=6)
    style_body = ParagraphStyle('Body', parent=styles['Normal'], fontSize=10, textColor=color_text, leading=14)
    style_sidebar_item = ParagraphStyle('SidebarItem', parent=styles['Normal'], fontSize=9, textColor=color_text_light, leading=13)
    style_date = ParagraphStyle('Date', parent=styles['Normal'], fontSize=9, textColor=color_accent, fontName='Helvetica-Bold', alignment=TA_RIGHT)

    # --- Contenu Sidebar ---
    sidebar_content = []
    
    # Image de profil
    profile_pic = cv_data.get('profile_picture')
    if profile_pic:
        try:
            from reportlab.platypus import Image
            import os
            
            # Déterminer le chemin local ou URL
            img_path = None
            if profile_pic.startswith('/static/'):
                img_path = os.path.join(os.path.dirname(__file__), profile_pic.lstrip('/'))
            elif profile_pic.startswith('http'):
                # Pour les URLs (Google Profile), on pourrait les télécharger
                # Mais pour l'instant, on ignore ou on utilise un placeholder si distant
                pass
                
            if img_path and os.path.exists(img_path):
                img = Image(img_path, width=45*mm, height=45*mm)
                sidebar_content.append(img)
                sidebar_content.append(Spacer(1, 10*mm))
        except Exception as e:
            print(f"Erreur image PDF: {e}")

    sidebar_content.append(Paragraph("COORDONNÉES", style_sidebar_h))
    sidebar_content.append(Paragraph(f"Localisation: {cv_data.get('localisation', 'France')}", style_sidebar_item))
    sidebar_content.append(Paragraph(f"Email: {cv_data.get('email', '')}", style_sidebar_item))
    sidebar_content.append(Paragraph(f"Tel: {cv_data.get('telephone', '')}", style_sidebar_item))
    
    sidebar_content.append(Paragraph("COMPÉTENCES", style_sidebar_h))
    for comp in cv_data.get('competences', [])[:12]:
        sidebar_content.append(Paragraph(f"• {comp}", style_sidebar_item))
    
    sidebar_content.append(Paragraph("LANGUES", style_sidebar_h))
    for langue in cv_data.get('langues', ['Français : Courant']):
        sidebar_content.append(Paragraph(f"• {langue}", style_sidebar_item))

    # --- Contenu Principal ---
    main_content = []
    # Header
    main_content.append(Paragraph(f"{cv_data.get('prenom', '')} {cv_data.get('nom', '')}".upper(), style_name))
    main_content.append(Paragraph(cv_data.get('titre_professionnel', 'PROFESSIONNEL'), style_title))
    main_content.append(Spacer(1, 10*mm))
    
    # Expériences
    main_content.append(Paragraph("EXPÉRIENCES PROFESSIONNELLES", style_section))
    for exp in cv_data.get('experiences', []):
        # Sous-table pour le titre et la date
        exp_header = Table([
            [Paragraph(f"<b>{exp.get('titre', '')}</b>", style_body), Paragraph(exp.get('periode', ''), style_date)]
        ], colWidths=[100*mm, 35*mm])
        exp_header.setStyle(TableStyle([('LEFTPADDING', (0,0), (-1,-1), 0), ('RIGHTPADDING', (0,0), (-1,-1), 0)]))
        main_content.append(exp_header)
        
        main_content.append(Paragraph(f"<i>{exp.get('entreprise', '')}</i>", style_body))
        for tache in exp.get('taches', []):
            main_content.append(Paragraph(f"• {tache}", style_body))
        main_content.append(Spacer(1, 4*mm))
    
    # Formations
    main_content.append(Paragraph("FORMATIONS", style_section))
    for form in cv_data.get('formations', []):
        form_header = Table([
            [Paragraph(f"<b>{form.get('titre', '')}</b>", style_body), Paragraph(form.get('periode', ''), style_date)]
        ], colWidths=[100*mm, 35*mm])
        form_header.setStyle(TableStyle([('LEFTPADDING', (0,0), (-1,-1), 0), ('RIGHTPADDING', (0,0), (-1,-1), 0)]))
        main_content.append(form_header)
        main_content.append(Paragraph(form.get('etablissement', ''), style_body))
        main_content.append(Spacer(1, 3*mm))

    # Layout final : Table de 2 colonnes
    # Sidebar (65mm) | Main (135mm) = 200mm (A4 est 210mm)
    data = [[sidebar_content, main_content]]
    layout_table = Table(data, colWidths=[65*mm, 135*mm], rowHeights=None) 
    layout_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), color_sidebar_bg),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (0,0), 8*mm),
        ('RIGHTPADDING', (0,0), (0,0), 8*mm),
        ('TOPPADDING', (0,0), (-1,-1), 12*mm),
        ('BOTTOMPADDING', (0,0), (-1,-1), 12*mm),
        ('LEFTPADDING', (1,0), (1,0), 10*mm),
        ('RIGHTPADDING', (1,0), (1,0), 10*mm),
    ]))

    story = [layout_table]
    doc.build(story)
    buffer.seek(0)
    return buffer

@app.route('/api/download-cv-pdf', methods=['POST'])
def download_cv_pdf():
    """Télécharge le CV amélioré en PDF"""
    try:
        data = request.get_json()
        cv_data = data.get('cv_data')
        nom = data.get('nom', 'Mon_CV')
        
        if not cv_data:
            return jsonify({'error': 'Données CV manquantes'}), 400
        
        buffer = generate_cv_pdf_buffer(cv_data)
        
        return send_file(
            buffer,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=f'{nom}.pdf'
        )
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/send-cv-email', methods=['POST'])
@login_required
def send_cv_email_route():
    """Envoie le CV optimisé par email en PDF"""
    try:
        data = request.get_json()
        cv_data = data.get('cv_data')
        email = data.get('email', '').strip()
        nom = data.get('nom', 'CV_Optimise')
        
        if not cv_data:
            return jsonify({'error': 'Données CV manquantes'}), 400
        
        if not email:
            return jsonify({'error': 'Adresse email manquante'}), 400
        
        # Générer le PDF
        pdf_buffer = generate_cv_pdf_buffer(cv_data)
        filename = f"{nom}.pdf"
        
        # Envoyer par email
        from emails import send_cv_email
        success = send_cv_email(email, pdf_buffer, filename)
        
        if success:
            return jsonify({
                'success': True,
                'message': f'CV envoyé avec succès à {email}'
            })
        else:
            return jsonify({'error': 'Échec de l\'envoi de l\'email. Vérifiez la configuration SMTP.'}), 500
    
    except Exception as e:
        logger.error(f"Erreur envoi CV par email: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/download-cv-pptx', methods=['POST'])
def download_cv_pptx():
    """Télécharge le CV amélioré en PPTX"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        data = request.get_json()
        cv_data = data.get('cv_data')
        nom = data.get('nom', 'Mon_CV')
        
        if not cv_data:
            return jsonify({'error': 'Données CV manquantes'}), 400
        
        # Créer la présentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Titre et Contact
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Titre
        title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = f"{cv_data.get('prenom', '')} {cv_data.get('nom', '')}".upper()
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.add_paragraph()
        subtitle_para.text = cv_data.get('titre_professionnel', 'PROFESSIONNEL')
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(52, 152, 219)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Contact
        contact_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.add_paragraph()
        contact_para.text = f"{cv_data.get('email', '')} | {cv_data.get('telephone', '')} | {cv_data.get('localisation', '')}"
        contact_para.font.size = Pt(14)
        contact_para.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Expériences
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        
        exp_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        exp_title_frame = exp_title.text_frame
        exp_title_para = exp_title_frame.add_paragraph()
        exp_title_para.text = "EXPÉRIENCES PROFESSIONNELLES"
        exp_title_para.font.size = Pt(24)
        exp_title_para.font.bold = True
        exp_title_para.font.color.rgb = RGBColor(52, 152, 219)
        
        y_pos = 1.3
        for exp in cv_data.get('experiences', [])[:3]:
            exp_box = slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
            exp_frame = exp_box.text_frame
            
            exp_p1 = exp_frame.add_paragraph()
            exp_p1.text = f"{exp.get('titre', '')} - {exp.get('periode', '')}"
            exp_p1.font.size = Pt(14)
            exp_p1.font.bold = True
            
            exp_p2 = exp_frame.add_paragraph()
            exp_p2.text = exp.get('entreprise', '')
            exp_p2.font.size = Pt(12)
            
            for tache in exp.get('taches', [])[:2]:
                exp_p3 = exp_frame.add_paragraph()
                exp_p3.text = f"• {tache[:80]}"
                exp_p3.font.size = Pt(10)
                exp_p3.level = 1
            
            y_pos += 1.8
        
        # Slide 3: Compétences
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        
        comp_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        comp_title_frame = comp_title.text_frame
        comp_title_para = comp_title_frame.add_paragraph()
        comp_title_para.text = "COMPÉTENCES & FORMATIONS"
        comp_title_para.font.size = Pt(24)
        comp_title_para.font.bold = True
        comp_title_para.font.color.rgb = RGBColor(52, 152, 219)
        
        # Compétences
        comp_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(2))
        comp_frame = comp_box.text_frame
        comp_p = comp_frame.add_paragraph()
        comp_p.text = "Compétences:"
        comp_p.font.size = Pt(16)
        comp_p.font.bold = True
        
        for comp in cv_data.get('competences', [])[:8]:
            comp_item = comp_frame.add_paragraph()
            comp_item.text = f"• {comp}"
            comp_item.font.size = Pt(12)
            comp_item.level = 1
        
        # Formations
        form_p = comp_frame.add_paragraph()
        form_p.text = "\nFormations:"
        form_p.font.size = Pt(16)
        form_p.font.bold = True
        
        for form in cv_data.get('formations', []):
            form_item = comp_frame.add_paragraph()
            form_item.text = f"• {form.get('titre', '')} - {form.get('etablissement', '')}"
            form_item.font.size = Pt(12)
            form_item.level = 1
        
        # Sauvegarder en mémoire
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return send_file(
            buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'{nom}.pptx'
        )
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============= GEMINI ENDPOINTS SUPPLÉMENTAIRES =============

@app.route('/api/enhance-cv-gemini', methods=['POST'])
def enhance_cv_gemini():
    """Améliore un CV avec les suggestions de Gemini"""
    try:
        if not gemini_service.enabled:
            return jsonify({'error': 'Service Gemini non disponible'}), 503
        
        data = request.get_json()
        cv_text = data.get('cv_text', '').strip()
        job_title = data.get('job_title', '').strip()
        
        if not cv_text or len(cv_text) < 100:
            return jsonify({'error': 'CV invalide (minimum 100 caractères)'}), 400
        
        # Obtenir les suggestions avec Gemini
        suggestions = gemini_service.enhance_cv(
            cv_text=cv_text,
            job_title=job_title if job_title else None
        )
        
        if 'error' in suggestions:
            return jsonify({'error': 'Erreur lors de l\'amélioration du CV', 'details': suggestions.get('error')}), 500
        
        return jsonify({
            'success': True,
            'suggestions': suggestions,
            'generated_at': datetime.utcnow().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/optimize-cv-for-offer', methods=['POST'])
def optimize_cv_for_offer():
    """Optimise le CV en fonction d'une offre d'emploi spécifique"""
    try:
        if not gemini_service.enabled:
            return jsonify({'error': 'Service IA non disponible'}), 503
        
        data = request.get_json()
        cv_text = data.get('cv_text', '').strip()
        job_offer_text = data.get('job_offer_text', '').strip()
        
        if not cv_text or len(cv_text) < 100:
            return jsonify({'error': 'CV invalide (minimum 100 caractères)'}), 400
        
        if not job_offer_text or len(job_offer_text) < 50:
            return jsonify({'error': 'Offre d\'emploi invalide'}), 400
        
        # Optimiser le CV en fonction de l'offre
        result = gemini_service.optimize_cv_for_job_offer(
            cv_text=cv_text,
            job_offer_text=job_offer_text
        )
        
        if 'error' in result:
            return jsonify({'error': result.get('error')}), 500
        
        return jsonify({
            'success': True,
            'optimization': result,
            'generated_at': datetime.utcnow().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyze-compatibility-ai', methods=['POST'])
def analyze_compatibility_ai():
    """Analyse avancée de compatibilité avec Gemini (sans PDF)"""
    try:
        if not gemini_service.enabled:
            return jsonify({'error': 'Service Gemini non disponible'}), 503
        
        data = request.get_json()
        cv_text = data.get('cv_text', '').strip()
        job_offer = data.get('job_offer', '').strip()
        
        if not cv_text or len(cv_text) < 50:
            return jsonify({'error': 'CV invalide (minimum 50 caractères)'}), 400
        
        if not job_offer or len(job_offer) < 50:
            return jsonify({'error': 'Offre invalide (minimum 50 caractères)'}), 400
        
        # Analyse avec Gemini
        analysis = gemini_service.analyze_cv_compatibility(cv_text, job_offer)
        
        if 'error' in analysis:
            return jsonify({'error': 'Erreur lors de l\'analyse', 'details': analysis.get('error')}), 500
        
        # Optionnel : sauvegarder si utilisateur connecté
        if current_user.is_authenticated:
            try:
                job_title = job_offer[:100] + "..." if len(job_offer) > 100 else job_offer
                cv_analysis = CVAnalysis(
                    user_id=current_user.id,
                    filename='Text Analysis',
                    job_title=job_title,
                    final_score=analysis.get('compatibility_score', 0),
                    similarity_score=0,
                    coverage_score=0,
                    recommendation=analysis.get('insight', 'Analyse effectuée'),
                    missing_keywords=json.dumps(analysis.get('missing_skills', []))
                )
                db.session.add(cv_analysis)
                db.session.commit()
            except Exception as e:
                print(f"Erreur sauvegarde: {e}")
                db.session.rollback()
        
        return jsonify({
            'success': True,
            'analysis': analysis,
            'generated_at': datetime.utcnow().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/chat', methods=['POST'])
def api_chat():
    """Endpoint pour le chatbot Gemini"""
    try:
        data = request.get_json()
        message = data.get('message', '').strip()
        history = data.get('history', [])

        if not message:
            return jsonify({'error': 'Message vide'}), 400

        # Appel au service Gemini
        response = gemini_service.chat(message, history)

        return jsonify({
            'success': True,
            'response': response
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ============= CANDIDATURES & OFFRES SAUVEGARDÉES =============

@app.route('/api/save-offer', methods=['POST'])
@login_required
def save_offer():
    """Sauvegarde une offre d'emploi"""
    try:
        data = request.get_json()
        
        if not data.get('title') or not data.get('company'):
            return jsonify({'error': 'Données incomplètes'}), 400
        
        from models import JobOffer, SavedOffer
        
        # Créer ou trouver l'offre
        offer = JobOffer(
            title=data.get('title'),
            company=data.get('company'),
            location=data.get('location', ''),
            salary=data.get('salary', ''),
            description=data.get('description', ''),
            contract_type=data.get('contract_type', 'CDI')
        )
        db.session.add(offer)
        db.session.flush()
        
        # Sauvegarder l'offre pour l'utilisateur
        saved = SavedOffer(
            user_id=current_user.id,
            offer_id=offer.id,
            status='saved'
        )
        db.session.add(saved)
        db.session.commit()
        
        logger.info(f"Offre sauvegardée pour user {current_user.id}: {offer.title}")
        
        return jsonify({
            'success': True,
            'message': 'Offre sauvegardée avec succès',
            'offer_id': offer.id
        })
        
    except Exception as e:
        db.session.rollback()
        logger.error(f"Erreur sauvegarde offre: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/get-applications', methods=['GET'])
@login_required
def get_applications():
    """Récupère toutes les candidatures de l'utilisateur"""
    try:
        from models import SavedOffer
        
        applications = SavedOffer.query.filter_by(user_id=current_user.id).order_by(SavedOffer.created_at.desc()).all()
        
        app_list = []
        for app in applications:
            if app.offer:
                app_list.append({
                    'id': app.id,
                    'title': app.offer.title,
                    'company': app.offer.company,
                    'location': app.offer.location,
                    'salary': app.offer.salary,
                    'description': app.offer.description,
                    'contract_type': app.offer.contract_type,
                    'status': app.status,
                    'date': app.created_at.strftime('%d %b %Y'),
                    'offer_id': app.offer.id
                })
        
        return jsonify({
            'success': True,
            'applications': app_list,
            'count': len(app_list)
        })
        
    except Exception as e:
        logger.error(f"Erreur récupération candidatures: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/update-application-status', methods=['POST'])
@login_required
def update_application_status():
    """Mets à jour le statut d'une candidature"""
    try:
        data = request.get_json()
        application_id = data.get('id')
        new_status = data.get('status')  # 'saved', 'applied', 'rejected'
        
        if not application_id or not new_status:
            return jsonify({'error': 'Données incomplètes'}), 400
        
        from models import SavedOffer
        
        app = SavedOffer.query.filter_by(id=application_id, user_id=current_user.id).first()
        if not app:
            return jsonify({'error': 'Candidature non trouvée'}), 404
        
        app.status = new_status
        db.session.commit()
        
        logger.info(f"Statut candidature {application_id} mis à jour: {new_status}")
        
        return jsonify({
            'success': True,
            'message': f'Statut mis à jour en "{new_status}"',
            'status': new_status
        })
        
    except Exception as e:
        db.session.rollback()
        logger.error(f"Erreur mise à jour statut: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/delete-application', methods=['POST'])
@login_required
def delete_application():
    """Supprime une candidature"""
    try:
        data = request.get_json()
        application_id = data.get('id')
        
        if not application_id:
            return jsonify({'error': 'ID manquant'}), 400
        
        from models import SavedOffer
        
        app = SavedOffer.query.filter_by(id=application_id, user_id=current_user.id).first()
        if not app:
            return jsonify({'error': 'Candidature non trouvée'}), 404
        
        db.session.delete(app)
        db.session.commit()
        
        logger.info(f"Candidature {application_id} supprimée")
        
        return jsonify({
            'success': True,
            'message': 'Candidature supprimée'
        })
        
    except Exception as e:
        db.session.rollback()
        logger.error(f"Erreur suppression candidature: {str(e)}")
        return jsonify({'error': str(e)}), 500


# ============= USER MANAGEMENT API ENDPOINTS =============

@app.route('/api/users/<int:user_id>/update', methods=['POST'])
@login_required
def update_user(user_id):
    """Update user details (admin only)"""
    try:
        # Check if current user is admin
        if not current_user.is_admin:
            return jsonify({'success': False, 'error': 'Accès non autorisé'}), 403
        
        # Get the user to update
        user = User.query.get(user_id)
        if not user:
            return jsonify({'success': False, 'error': 'Utilisateur introuvable'}), 404
        
        # Prevent self-demotion
        if user.id == current_user.id:
            return jsonify({'success': False, 'error': 'Vous ne pouvez pas modifier votre propre rôle'}), 400
        
        # Get update data
        data = request.get_json()
        is_admin = data.get('is_admin', False)
        
        # Update user role
        user.is_admin = is_admin
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': 'Utilisateur mis à jour avec succès',
            'user': {
                'id': user.id,
                'is_admin': user.is_admin
            }
        })
        
    except Exception as e:
        db.session.rollback()
        print(f"Error updating user: {e}")
        return jsonify({'success': False, 'error': 'Erreur lors de la mise à jour'}), 500


@app.route('/api/users/<int:user_id>/delete', methods=['DELETE'])
@login_required
def delete_user_admin(user_id):
    """Delete a user (admin only)"""
    try:
        if not current_user.is_admin:
            return jsonify({'success': False, 'error': 'Accès non autorisé'}), 403
        
        user = User.query.get(user_id)
        if not user:
            return jsonify({'success': False, 'error': 'Utilisateur introuvable'}), 404
        
        if user.id == current_user.id:
            return jsonify({'success': False, 'error': 'Vous ne pouvez pas supprimer votre propre compte'}), 400
        
        user_name = f"{user.firstname} {user.lastname}"
        db.session.delete(user)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'Utilisateur {user_name} supprimé avec succès'
        })
    except Exception as e:
        db.session.rollback()
        print(f"Error deleting user: {e}")
        return jsonify({'success': False, 'error': 'Erreur lors de la suppression'}), 500


@app.route('/api/chatbot', methods=['POST'])
@login_required
def chatbot_message():
    """Handle chatbot messages with Gemini AI"""
    try:
        data = request.get_json()
        message = data.get('message', '').strip()
        
        if not message:
            return jsonify({'error': 'Message vide'}), 400

        # Import ChatMessage model locally to avoid circular imports if any
        from models import ChatMessage
        
        # 1. Save User Message
        user_msg = ChatMessage(
            user_id=current_user.id,
            session_id=str(current_user.id), # Use user ID as session ID for logged users
            role='user',
            message=message
        )
        db.session.add(user_msg)
        db.session.commit()
        
        # Add context about JobMatch platform
        context_message = f"""Tu es l'assistant IA de JobMatch, une plateforme d'optimisation de CV et d'analyse d'offres d'emploi. 
        L'utilisateur {current_user.firstname} {current_user.lastname} te pose la question suivante: {message}
        
        Réponds de manière concise et professionnelle en français. Si la question concerne:
        - L'analyse de CV: explique comment utiliser notre outil d'analyse
        - Les offres d'emploi: guide vers l'analyseur d'offres
        - Le profil: explique comment mettre à jour les informations
        - Les tendances: parle des tendances du marché de l'emploi
        
        Garde ta réponse courte (maximum 3-4 phrases)."""
        
        # Get response from Gemini
        response_text = gemini_service.chat(context_message)
        
        # 2. Save Bot Response
        bot_msg = ChatMessage(
            user_id=current_user.id,
            session_id=str(current_user.id),
            role='bot',
            message=response_text
        )
        db.session.add(bot_msg)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'response': response_text
        })
        
    except Exception as e:
        logger.error(f"Erreur chatbot: {str(e)}")
        # Rollback in case of DB error
        db.session.rollback()
        return jsonify({
            'success': False,
            'error': 'Erreur lors du traitement de votre message',
            'response': "Désolé, je rencontre une difficulté technique. Veuillez réessayer."
        }), 500

@app.route('/api/chatbot/history', methods=['GET'])
@login_required
def get_chat_history():
    """Get chat history for the current user"""
    try:
        from models import ChatMessage
        
        # Retrieve last 50 messages ordered by time
        messages = ChatMessage.query.filter_by(user_id=current_user.id)\
            .order_by(ChatMessage.created_at.asc())\
            .limit(50)\
            .all()
            
        history = [{
            'role': msg.role,
            'message': msg.message,
            'created_at': msg.created_at.isoformat()
        } for msg in messages]
        
        return jsonify({'success': True, 'history': history})
    except Exception as e:
        logger.error(f"Erreur historique chatbot: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

# ============= ROUTES GESTION DES CANDIDATURES (CANDIDACY) =============

@app.route('/api/candidacies', methods=['GET'])
@login_required
def list_candidacies():
    """
    Retourne toutes les candidatures de l'utilisateur.
    Créées automatiquement à chaque analyse ou optimisation de CV.
    """
    try:
        candidacies = Candidacy.query.filter_by(user_id=current_user.id)\
            .order_by(Candidacy.created_at.desc()).all()
        
        result = []
        for c in candidacies:
            result.append({
                'id': c.id,
                'job_title': c.job_title,
                'company': c.company or '',
                'filename': c.filename or '',
                'cv_score': c.cv_score,
                'source': c.source,  # 'analyzed', 'optimized', 'manual'
                'status': c.status,  # 'pending', 'submitted', 'won', 'lost'
                'is_submitted': c.is_submitted,
                'feedback': c.feedback or '',
                'created_at': c.created_at.strftime('%d/%m/%Y'),
                'submitted_at': c.submitted_at.strftime('%d/%m/%Y') if c.submitted_at else None,
                'resolved_at': c.resolved_at.strftime('%d/%m/%Y') if c.resolved_at else None,
            })
        
        return jsonify({'success': True, 'candidacies': result})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/candidacy/add', methods=['POST'])
@login_required
def add_candidacy():
    """Ajout manuel d'une candidature."""
    try:
        data = request.get_json()
        if not data or not data.get('job_title'):
            return jsonify({'error': 'Le titre du poste est requis'}), 400
        
        cand = Candidacy(
            user_id=current_user.id,
            job_title=data.get('job_title'),
            company=data.get('company', ''),
            source='manual',
            status='pending',
            is_submitted=False
        )
        db.session.add(cand)
        db.session.commit()
        
        return jsonify({'success': True, 'candidacy_id': cand.id})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/candidacy/submit/<int:id>', methods=['POST'])
@login_required
def submit_candidacy(id):
    """
    Marque un CV comme 'soumis' à une entreprise.
    L'utilisateur indique l'entreprise à ce moment-là.
    """
    try:
        c = Candidacy.query.get_or_404(id)
        if c.user_id != current_user.id:
            return jsonify({'error': 'Accès refusé'}), 403
        
        data = request.get_json()
        company = data.get('company', '').strip()
        
        if not company:
            return jsonify({'error': 'Nom de l\'entreprise requis'}), 400
        
        c.company = company
        c.is_submitted = True
        c.status = 'submitted'
        c.submitted_at = datetime.utcnow()
        db.session.commit()
        
        return jsonify({'success': True, 'message': f'CV soumis à {company} !'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/candidacy/won/<int:id>', methods=['POST'])
@login_required
def candidacy_won(id):
    """Marque la candidature comme gagnée (embauché)."""
    try:
        c = Candidacy.query.get_or_404(id)
        if c.user_id != current_user.id:
            return jsonify({'error': 'Accès refusé'}), 403
        
        data = request.get_json() or {}
        c.status = 'won'
        c.feedback = data.get('feedback', c.feedback)
        c.resolved_at = datetime.utcnow()
        db.session.commit()
        
        # Notification de succès
        try:
            notif = Notification(
                user_id=current_user.id,
                title='🎉 Félicitations !',
                message=f'Vous avez décroché le poste de {c.job_title} chez {c.company or "cette entreprise"} !',
                type='success'
            )
            db.session.add(notif)
            db.session.commit()
        except: pass
        
        return jsonify({'success': True, 'message': '🎉 Bravo ! Candidature gagnée enregistrée !'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/candidacy/lost/<int:id>', methods=['POST'])
@login_required
def candidacy_lost(id):
    """Marque la candidature comme perdue (refusé)."""
    try:
        c = Candidacy.query.get_or_404(id)
        if c.user_id != current_user.id:
            return jsonify({'error': 'Accès refusé'}), 403
        
        data = request.get_json() or {}
        c.status = 'lost'
        c.feedback = data.get('feedback', c.feedback)
        c.resolved_at = datetime.utcnow()
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Résultat enregistré. Ne vous découragez pas !'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/candidacy/delete/<int:id>', methods=['POST'])
@login_required
def remove_candidacy(id):
    """Supprime une candidature."""
    try:
        c = Candidacy.query.get_or_404(id)
        if c.user_id != current_user.id:
            return jsonify({'error': 'Accès refusé'}), 403
        db.session.delete(c)
        db.session.commit()
        return jsonify({'success': True, 'message': 'Candidature supprimée'})
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/stats/candidacies', methods=['GET'])
@login_required
def get_candidacy_stats():
    """Statistiques des candidatures (soumises uniquement pour le taux de succès)."""
    try:
        if current_user.is_admin:
            base = Candidacy.query
        else:
            base = Candidacy.query.filter_by(user_id=current_user.id)

        total = base.count()
        submitted = base.filter(Candidacy.is_submitted == True).count() if not current_user.is_admin \
            else Candidacy.query.filter(Candidacy.is_submitted == True).count()
        won = base.filter_by(status='won').count() if not current_user.is_admin \
            else Candidacy.query.filter_by(status='won').count()
        lost = base.filter_by(status='lost').count() if not current_user.is_admin \
            else Candidacy.query.filter_by(status='lost').count()
        pending = base.filter_by(status='pending').count() if not current_user.is_admin \
            else Candidacy.query.filter_by(status='pending').count()

        # Taux de succès = gagnées / soumises
        success_rate = round((won / submitted * 100), 1) if submitted > 0 else 0

        return jsonify({
            'success': True,
            'stats': {
                'total': total,
                'submitted': submitted,
                'won': won,
                'lost': lost,
                'pending': pending,
                'success_rate': success_rate
            }
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    # Utilisation du port assigné par l'environnement ou 5000 par défaut
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
